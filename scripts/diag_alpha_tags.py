

import argparse
from pptx import Presentation
from pptx.oxml.ns import qn


def find_first(el, tag):
    try:
        return el.find(qn(tag))
    except Exception:
        return None


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--pptx", required=True)
    ap.add_argument("--max_hits", type=int, default=80)
    args = ap.parse_args()

    prs = Presentation(args.pptx)

    targets = ("a:alpha", "a:alphaMod", "a:alphaOff")

    hits = []
    for si, slide in enumerate(prs.slides, start=1):
        for idx, shp in enumerate(slide.shapes, start=1):
            try:
                spPr = shp._element.spPr
            except Exception:
                continue

            found = []

            eff = find_first(spPr, "a:effectLst")
            if eff is not None:
                amf = find_first(eff, "a:alphaModFix")
                if amf is not None and amf.get("amt") is not None:
                    found.append(("effect.alphaModFix", amf.get("amt")))

            solid = find_first(spPr, "a:solidFill")
            if solid is not None:
                for tag in ("a:srgbClr", "a:schemeClr", "a:prstClr"):
                    c = find_first(solid, tag)
                    if c is None:
                        continue
                    for atag in targets:
                        ael = find_first(c, atag)
                        if ael is not None and ael.get("val") is not None:
                            found.append((f"solidFill.{tag}.{atag}", ael.get("val")))

            ln = find_first(spPr, "a:ln")
            if ln is not None:
                lsolid = find_first(ln, "a:solidFill")
                if lsolid is not None:
                    for tag in ("a:srgbClr", "a:schemeClr", "a:prstClr"):
                        c = find_first(lsolid, tag)
                        if c is None:
                            continue
                        for atag in targets:
                            ael = find_first(c, atag)
                            if ael is not None and ael.get("val") is not None:
                                found.append((f"line.{tag}.{atag}", ael.get("val")))

            if found:
                hits.append((si, idx, str(shp.shape_type), found))

    print("hits", len(hits))
    for h in hits[: args.max_hits]:
        si, idx, st, found = h
        print(si, idx, st, found)


if __name__ == "__main__":
    main()
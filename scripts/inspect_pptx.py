from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

path = "input/テスト2.pptx"
prs = Presentation(path)

total_text_shapes = 0
total_picture_shapes = 0
total_other_shapes = 0
slides_summary = []

for si, slide in enumerate(prs.slides, start=1):
    text_shapes = 0
    picture_shapes = 0
    other_shapes = 0
    text_chars = 0

    for shp in slide.shapes:
        st = shp.shape_type
        if st == MSO_SHAPE_TYPE.PICTURE:
            picture_shapes += 1
        elif getattr(shp, "has_text_frame", False) and shp.has_text_frame:
            text_shapes += 1
            txt = shp.text_frame.text or ""
            text_chars += len(txt.strip())
        else:
            other_shapes += 1

    total_text_shapes += text_shapes
    total_picture_shapes += picture_shapes
    total_other_shapes += other_shapes
    slides_summary.append((si, text_shapes, picture_shapes, other_shapes, text_chars))

print("slides:", len(prs.slides))
print("TOTAL text_shapes:", total_text_shapes)
print("TOTAL picture_shapes:", total_picture_shapes)
print("TOTAL other_shapes:", total_other_shapes)

print("\nTop 10 slides by text_chars:")
for si, ts, ps, os, tc in sorted(slides_summary, key=lambda x: x[4], reverse=True)[:10]:
    print(f"  slide {si:>3}: text_shapes={ts:>2}, pictures={ps:>2}, other={os:>2}, text_chars={tc}")

from __future__ import annotations

import hashlib
import re
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.dml.color import ColorFormat
from pptx.enum.dml import MSO_FILL

from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.oxml.ns import qn



def _slugify_ascii(name: str) -> str:
    s = name.strip()
    s = re.sub(r"\s+", "_", s)
    s = re.sub(r"[^a-zA-Z0-9_\-]", "_", s)
    s = re.sub(r"_+", "_", s).strip("_")
    if not s:
        s = "document"
    return s[:64]


def _norm_text(s: str) -> str:
    s = s.replace("\u00a0", " ")
    s = re.sub(r"[ \t]+", " ", s)
    s = re.sub(r"\n{3,}", "\n\n", s)
    return s.strip()


def _emu(v: Any) -> int:
    """python-pptx uses EMU for positions/sizes. Keep as int and clamp."""
    try:
        iv = int(v)
    except Exception:
        return 0
    return iv if iv >= 0 else 0


# Theme scheme keys used in ppt/theme/theme1.xml
_THEME_KEY_BY_ENUM_NAME: Dict[str, str] = {
    "DARK_1": "dk1",
    "LIGHT_1": "lt1",
    "DARK_2": "dk2",
    "LIGHT_2": "lt2",
    "ACCENT_1": "accent1",
    "ACCENT_2": "accent2",
    "ACCENT_3": "accent3",
    "ACCENT_4": "accent4",
    "ACCENT_5": "accent5",
    "ACCENT_6": "accent6",
    "HYPERLINK": "hlink",
    "FOLLOWED_HYPERLINK": "folHlink",
}


def _load_theme_rgb_map(pptx_path: Path) -> Dict[str, str]:
    """Load theme color scheme (best-effort) from ppt/theme/theme1.xml.

    Returns mapping like {'accent1': 'RRGGBB', 'dk1': 'RRGGBB', ...}
    """
    out: Dict[str, str] = {}
    try:
        with zipfile.ZipFile(pptx_path, "r") as zf:
            # theme file name is usually fixed
            name = "ppt/theme/theme1.xml"
            if name not in set(zf.namelist()):
                # fallback: any ppt/theme/theme*.xml
                themes = [n for n in zf.namelist() if n.startswith("ppt/theme/") and n.endswith(".xml")]
                if not themes:
                    return out
                name = sorted(themes)[0]
            xml_bytes = zf.read(name)
    except Exception:
        return out

    try:
        root = ET.fromstring(xml_bytes)
    except Exception:
        return out

    ns = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    }

    # a:theme/a:themeElements/a:clrScheme
    clr = root.find(".//a:themeElements/a:clrScheme", ns)
    if clr is None:
        return out

    for child in list(clr):
        # child tag ends with scheme key (dk1, lt1, accent1...)
        tag = child.tag
        if "}" in tag:
            key = tag.split("}", 1)[1]
        else:
            key = tag

        # Prefer srgbClr@val, else sysClr@lastClr
        srgb = child.find(".//a:srgbClr", ns)
        if srgb is not None and srgb.get("val"):
            out[key] = str(srgb.get("val")).upper()
            continue
        sysc = child.find(".//a:sysClr", ns)
        if sysc is not None and sysc.get("lastClr"):
            out[key] = str(sysc.get("lastClr")).upper()
            continue

    return out


def _rgb_from_color(color: Optional[ColorFormat], theme_rgb: Optional[Dict[str, str]] = None) -> Optional[str]:
    """Return hex RGB like 'RRGGBB' when available.

    If the color is a theme color, attempt to resolve via `theme_rgb`.
    """
    if color is None:
        return None

    # Direct RGB
    try:
        if getattr(color, "rgb", None) is not None:
            return str(color.rgb)
    except Exception:
        pass

    # Theme color -> scheme key -> RGB
    try:
        tc = getattr(color, "theme_color", None)
        if tc is not None and theme_rgb:
            enum_name = getattr(tc, "name", None) or str(tc)
            enum_name = str(enum_name)
            scheme_key = _THEME_KEY_BY_ENUM_NAME.get(enum_name)
            if scheme_key and scheme_key in theme_rgb:
                return theme_rgb[scheme_key]
    except Exception:
        pass

    return None


# Alpha helpers for extracting opacity from DrawingML XML (theme/preset colors)
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
# Alpha helpers for extracting opacity from DrawingML XML (theme/preset colors)

def _alpha_from_val_100000(val: Any) -> Optional[float]:
    try:
        iv = int(str(val))
    except Exception:
        return None
    if iv < 0:
        iv = 0
    if iv > 100000:
        iv = 100000
    return iv / 100000.0


def _effective_alpha_from_color_node(clr: Any) -> Optional[float]:
    """Return effective alpha from a color node.

    Supports:
    - a:alpha (absolute)
    - a:alphaMod / a:alphaOff (modifiers). For modifiers, assume base alpha=1.0:
      alpha = clamp(1.0 * alphaMod + alphaOff)
    Values are 0..100000.
    """
    if clr is None:
        return None

    # Absolute alpha wins
    try:
        ael = clr.find(qn("a:alpha"))
        if ael is not None and ael.get("val") is not None:
            return _alpha_from_val_100000(ael.get("val"))
    except Exception:
        pass

    mod = None
    off = None
    try:
        mel = clr.find(qn("a:alphaMod"))
        if mel is not None and mel.get("val") is not None:
            mod = _alpha_from_val_100000(mel.get("val"))
    except Exception:
        mod = None

    try:
        oel = clr.find(qn("a:alphaOff"))
        if oel is not None and oel.get("val") is not None:
            off = _alpha_from_val_100000(oel.get("val"))
    except Exception:
        off = None

    if mod is None and off is None:
        return None

    a = 1.0
    if mod is not None:
        a = a * mod
    if off is not None:
        a = a + off
    if a < 0.0:
        a = 0.0
    if a > 1.0:
        a = 1.0
    return a




def _xml_alpha_from_shape(shp: Any) -> Optional[float]:
    """Best-effort alpha (opacity) from DrawingML.

    We traverse via `find(qn(...))` to avoid xpath incompatibilities.

    Returns 0..1 where 1 is fully opaque.
    """
    try:
        spPr = shp._element.spPr
    except Exception:
        return None

    # Fill alpha: spPr/solidFill/(srgbClr|schemeClr|prstClr)/(alpha|alphaMod|alphaOff)
    try:
        solid = spPr.find(qn("a:solidFill"))
        if solid is not None:
            for tag in ("a:srgbClr", "a:schemeClr", "a:prstClr"):
                clr = solid.find(qn(tag))
                if clr is None:
                    continue
                a = _effective_alpha_from_color_node(clr)
                if a is not None:
                    return a
    except Exception:
        pass

    # Line alpha: spPr/ln/solidFill/(...)/(alpha|alphaMod|alphaOff)
    try:
        ln = spPr.find(qn("a:ln"))
        if ln is not None:
            lsolid = ln.find(qn("a:solidFill"))
            if lsolid is not None:
                for tag in ("a:srgbClr", "a:schemeClr", "a:prstClr"):
                    clr = lsolid.find(qn(tag))
                    if clr is None:
                        continue
                    a = _effective_alpha_from_color_node(clr)
                    if a is not None:
                        return a
    except Exception:
        pass

    return None




def _xml_alpha_from_line(shp: Any) -> Optional[float]:
    """Best-effort line alpha (opacity) from DrawingML line definition."""
    try:
        spPr = shp._element.spPr
    except Exception:
        return None

    try:
        ln = spPr.find(qn("a:ln"))
        if ln is None:
            return None
        solid = ln.find(qn("a:solidFill"))
        if solid is None:
            return None
        for tag in ("a:srgbClr", "a:schemeClr", "a:prstClr"):
            clr = solid.find(qn(tag))
            if clr is None:
                continue
            a = _effective_alpha_from_color_node(clr)
            if a is not None:
                return a
    except Exception:
        pass

    return None


def _xml_alpha_mod_fix_from_shape(shp: Any) -> Optional[float]:
    """Best-effort overall shape opacity from effect list.

    PowerPoint can encode object transparency as:
      a:spPr/a:effectLst/a:alphaModFix@amt
    where amt is 0..100000 (100000 = fully opaque).

    Returns 0..1 where 1 is fully opaque.
    """
    try:
        spPr = shp._element.spPr
    except Exception:
        return None

    try:
        eff = spPr.find(qn("a:effectLst"))
        if eff is None:
            return None
        amf = eff.find(qn("a:alphaModFix"))
        if amf is None:
            return None
        v = amf.get("amt")
        return _alpha_from_val_100000(v)
    except Exception:
        return None


def _xml_solid_color_rgb(shp: Any, theme_rgb: Optional[Dict[str, str]] = None) -> Optional[str]:
    """Best-effort solid fill RGB from underlying XML.

    Prefer a:solidFill/a:srgbClr@val. If schemeClr is used, resolve via theme_rgb.
    Returns hex RGB like 'RRGGBB'.
    """
    try:
        spPr = shp._element.spPr
    except Exception:
        return None

    try:
        solid = spPr.find(qn("a:solidFill"))
        if solid is None:
            return None

        srgb = solid.find(qn("a:srgbClr"))
        if srgb is not None and srgb.get("val"):
            return str(srgb.get("val")).upper()

        scheme = solid.find(qn("a:schemeClr"))
        if scheme is not None and scheme.get("val") and theme_rgb:
            key = str(scheme.get("val"))
            if key in theme_rgb:
                return str(theme_rgb[key]).upper()

        return None
    except Exception:
        return None


def _fill_dict(shp: Any, theme_rgb: Optional[Dict[str, str]] = None) -> Dict[str, Any]:
    out: Dict[str, Any] = {}
    try:
        fill = shp.fill
    except Exception:
        return out

    try:
        ftype = fill.type
    except Exception:
        ftype = None

    if ftype is None:
        return out

    # Only expose simple/robust fields (theme colors can be added later)
    if ftype == MSO_FILL.SOLID:
        out["type"] = "solid"
        try:
            rgb = _rgb_from_color(fill.fore_color, theme_rgb)
            if rgb is None:
                rgb = _xml_solid_color_rgb(shp, theme_rgb)
            out["color_rgb"] = rgb
        except Exception:
            pass
        try:
            # python-pptx transparency: 0..1 (1=fully transparent)
            tr = getattr(fill.fore_color, "transparency", None)
            if tr is not None:
                a = 1.0 - float(tr)
                if a < 0:
                    a = 0.0
                if a > 1:
                    a = 1.0
                out["alpha"] = a
        except Exception:
            pass
    elif ftype == MSO_FILL.BACKGROUND:
        out["type"] = "background"
    elif ftype == MSO_FILL.PATTERNED:
        out["type"] = "pattern"
    elif ftype == MSO_FILL.GRADIENT:
        out["type"] = "gradient"
    else:
        out["type"] = "other"

    # Fallbacks: transparency may be encoded only in the underlying XML, regardless of FillFormat.type
    if "alpha" not in out:
        aa = _xml_alpha_from_shape(shp)
        if aa is not None:
            out["alpha"] = aa

    if "alpha" not in out:
        aa = _xml_alpha_mod_fix_from_shape(shp)
        if aa is not None:
            out["alpha"] = aa

    return out

# Helper to extract fill info from a FillFormat (e.g., slide.background.fill)
def _fill_from_fillformat(fill: Any, theme_rgb: Optional[Dict[str, str]] = None) -> Dict[str, Any]:
    """Like _fill_dict but works on a FillFormat directly (e.g., slide.background.fill)."""
    out: Dict[str, Any] = {}
    try:
        ftype = fill.type
    except Exception:
        return out

    if ftype is None:
        return out

    if ftype == MSO_FILL.SOLID:
        out["type"] = "solid"
        try:
            out["color_rgb"] = _rgb_from_color(fill.fore_color, theme_rgb)
        except Exception:
            pass
        try:
            tr = getattr(fill.fore_color, "transparency", None)
            if tr is not None:
                a = 1.0 - float(tr)
                if a < 0:
                    a = 0.0
                if a > 1:
                    a = 1.0
                out["alpha"] = a
        except Exception:
            pass
        # Fallback: FillFormat does not expose underlying XML here; no fallback available.
    elif ftype == MSO_FILL.BACKGROUND:
        out["type"] = "background"
    elif ftype == MSO_FILL.PATTERNED:
        out["type"] = "pattern"
    elif ftype == MSO_FILL.GRADIENT:
        out["type"] = "gradient"
    else:
        out["type"] = "other"

    return out


def _line_dict(shp: Any, theme_rgb: Optional[Dict[str, str]] = None) -> Dict[str, Any]:
    out: Dict[str, Any] = {}
    try:
        line = shp.line
    except Exception:
        return out

    try:
        # width is EMU
        if line.width is not None:
            out["width_emu"] = _emu(line.width)
    except Exception:
        pass

    try:
        out["color_rgb"] = _rgb_from_color(line.color, theme_rgb)
    except Exception:
        pass
    try:
        tr = getattr(line.color, "transparency", None)
        if tr is not None:
            a = 1.0 - float(tr)
            if a < 0:
                a = 0.0
            if a > 1:
                a = 1.0
            out["alpha"] = a
    except Exception:
        pass
    # Fallback: some theme/preset line colors carry alpha only in the underlying XML
    if "alpha" not in out:
        aa = _xml_alpha_from_line(shp)
        if aa is not None:
            out["alpha"] = aa

    # Final fallback: overall shape opacity encoded as alphaModFix effect
    if "alpha" not in out:
        aa = _xml_alpha_mod_fix_from_shape(shp)
        if aa is not None:
            out["alpha"] = aa

    # If alpha is 0 (fully transparent), mark the line as invisible explicitly.
    # This ensures downstream renderers treat it as no-line even if they don't check alpha.
    if isinstance(out.get("alpha"), float) and out["alpha"] <= 0.0:
        out["visible"] = False

    return out


def _text_structure(shp: Any, theme_rgb: Optional[Dict[str, str]] = None) -> Dict[str, Any]:
    """Extract paragraph/run-level styling when available."""
    out: Dict[str, Any] = {"paragraphs": []}

    tf = getattr(shp, "text_frame", None)
    if tf is None:
        return out

    for pi, p in enumerate(tf.paragraphs):
        pinfo: Dict[str, Any] = {
            "index": pi,
            "alignment": str(getattr(p, "alignment", None)),
            "level": int(getattr(p, "level", 0) or 0),
            "runs": [],
        }

        for ri, r in enumerate(p.runs):
            txt = r.text or ""
            if not txt:
                continue

            font = getattr(r, "font", None)
            rinfo: Dict[str, Any] = {
                "index": ri,
                "text": txt,
            }

            if font is not None:
                try:
                    if font.name:
                        rinfo["font_name"] = font.name
                except Exception:
                    pass
                try:
                    if font.size is not None:
                        rinfo["font_size_emu"] = _emu(font.size)
                except Exception:
                    pass
                try:
                    if font.bold is not None:
                        rinfo["bold"] = bool(font.bold)
                except Exception:
                    pass
                try:
                    if font.italic is not None:
                        rinfo["italic"] = bool(font.italic)
                except Exception:
                    pass
                try:
                    if font.underline is not None:
                        rinfo["underline"] = bool(font.underline)
                except Exception:
                    pass
                try:
                    rinfo["color_rgb"] = _rgb_from_color(font.color, theme_rgb)
                except Exception:
                    pass
                try:
                    hc = getattr(font, "highlight_color", None)
                    hrgb = _rgb_from_color(hc, theme_rgb)
                    if hrgb:
                        rinfo["highlight_rgb"] = hrgb
                except Exception:
                    pass

            pinfo["runs"].append(rinfo)

        out["paragraphs"].append(pinfo)

    return out


def _shape_kind(shp: Any) -> str:
    """Coarse kind classification."""
    try:
        st = shp.shape_type
    except Exception:
        return "unknown"

    # Tables
    if getattr(shp, "has_table", False):
        return "table"

    # Pictures
    if st == MSO_SHAPE_TYPE.PICTURE:
        return "image"

    # Text (includes placeholders with text)
    if getattr(shp, "has_text_frame", False) and shp.has_text_frame:
        return "text"

    # Group
    if st == MSO_SHAPE_TYPE.GROUP:
        return "group"

    # Chart / SmartArt etc.
    if st == MSO_SHAPE_TYPE.CHART:
        return "chart"
    if st == MSO_SHAPE_TYPE.SMART_ART:
        return "smartart"

    return "shape"


# Map PowerPoint auto-shape types to renderer-friendly names.
def _autoshape_kind(shp: Any) -> Optional[str]:
    """Return a renderer-friendly shape kind for auto-shapes when possible."""
    try:
        st = shp.shape_type
    except Exception:
        return None

    if st == MSO_SHAPE_TYPE.LINE:
        return "line"

    if st != MSO_SHAPE_TYPE.AUTO_SHAPE:
        return None

    try:
        at = shp.auto_shape_type
    except Exception:
        return None

    # Common mappings used by the DOM renderer
    if at in (MSO_AUTO_SHAPE_TYPE.OVAL,):
        return "ellipse"
    if at in (MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGULAR_CALLOUT):
        return "round_rect"
    if at in (MSO_AUTO_SHAPE_TYPE.RECTANGLE, MSO_AUTO_SHAPE_TYPE.ROUND_1_RECTANGLE, MSO_AUTO_SHAPE_TYPE.ROUND_2_SAME_RECTANGLE, MSO_AUTO_SHAPE_TYPE.ROUND_2_DIAG_RECTANGLE):
        return "rect"

    # Fallback: still return rect so the renderer draws something rather than skipping
    return "rect"


def _image_info(shp: Any) -> Dict[str, Any]:
    """Return stable identifiers for embedded images. Writing files is handled elsewhere."""
    out: Dict[str, Any] = {}
    try:
        img = shp.image
    except Exception:
        return out

    try:
        blob = img.blob
        sha = hashlib.sha256(blob).hexdigest()
        out["sha256"] = sha
        out["ext"] = (img.ext or "").lstrip(".")
        out["byte_size"] = len(blob)
    except Exception:
        pass

    try:
        out["filename"] = img.filename
    except Exception:
        pass

    return out


def _table_info(shp: Any, theme_rgb: Optional[Dict[str, str]] = None) -> Dict[str, Any]:
    out: Dict[str, Any] = {"rows": 0, "cols": 0, "cells": []}
    try:
        tbl = shp.table
    except Exception:
        return out

    try:
        out["rows"] = int(len(tbl.rows))
        out["cols"] = int(len(tbl.columns))
    except Exception:
        pass

    for r in range(out["rows"]):
        for c in range(out["cols"]):
            cell = tbl.cell(r, c)
            text = _norm_text(getattr(cell, "text", "") or "")
            cell_info: Dict[str, Any] = {"r": r, "c": c, "text": text}
            # Capture per-run font info (font_size, font_name, bold, italic, color).
            # _text_structure works on any object with a text_frame attribute.
            try:
                ts = _text_structure(cell, theme_rgb)
                if ts.get("paragraphs"):
                    cell_info["paragraphs"] = ts["paragraphs"]
            except Exception:
                pass
            out["cells"].append(cell_info)

    return out


def extract_pptx(path: str | Path, *, include_extended: bool = False) -> dict[str, Any]:
    p = Path(path)
    theme_rgb = _load_theme_rgb_map(p)
    prs = Presentation(str(p))

    document_id = _slugify_ascii(p.stem)
    chunks: List[Dict[str, Any]] = []

    # Best-effort page size (EMU)
    page = {
        "w_emu": _emu(getattr(prs, "slide_width", 0)),
        "h_emu": _emu(getattr(prs, "slide_height", 0)),
    }

    for slide_idx, slide in enumerate(prs.slides, start=1):
        # Slide shapes are already in z-order (back-to-front) in most cases; keep index for stability.

        # Slide background (master/background fill). Many decks use this instead of a rectangle shape.
        if include_extended:
            try:
                bg_fill = _fill_from_fillformat(slide.background.fill, theme_rgb)
            except Exception:
                bg_fill = {}

            if bg_fill:
                bg_chunk: Dict[str, Any] = {
                    "chunk_id": f"s{slide_idx:03d}_bg000",
                    "page": slide_idx,
                    "bbox": {"x": 0, "y": 0, "w": page["w_emu"], "h": page["h_emu"]},
                    "text": "[BACKGROUND]",
                    "normalized_text": "[BACKGROUND]",
                    "heading_level": 0,
                    "kind": "shape",
                    "z": 0,
                    "rotation_deg": 0.0,
                    "style": {"fill": bg_fill, "line": {}},
                    "shape_kind": "rect",
                }
                chunks.append(bg_chunk)

        # Layout/Master shapes (background decorations) are not part of slide.shapes in many cases.
        # Extract them best-effort so backgrounds behind text are not lost.
        if include_extended:
            def _emit_from_shape_collection(shapes: Any, prefix: str, z_base: int) -> int:
                z = z_base
                try:
                    it = list(shapes)
                except Exception:
                    return z
                for i, shp2 in enumerate(it, start=1):
                    try:
                        kind2 = _shape_kind(shp2)
                    except Exception:
                        kind2 = "shape"

                    bbox2 = {
                        "x": _emu(getattr(shp2, "left", 0)),
                        "y": _emu(getattr(shp2, "top", 0)),
                        "w": _emu(getattr(shp2, "width", 0)),
                        "h": _emu(getattr(shp2, "height", 0)),
                    }

                    try:
                        rot2 = float(getattr(shp2, "rotation", 0) or 0)
                    except Exception:
                        rot2 = 0.0

                    base2: Dict[str, Any] = {
                        "chunk_id": f"s{slide_idx:03d}_{prefix}{i:03d}",
                        "page": slide_idx,
                        "bbox": bbox2,
                        "text": "",
                        "normalized_text": "",
                        "heading_level": 0,
                        "kind": kind2,
                        "z": z,
                        "rotation_deg": rot2,
                        "style": {
                            "fill": _fill_dict(shp2, theme_rgb),
                            "line": _line_dict(shp2, theme_rgb),
                        },
                    }

                    sk2 = _autoshape_kind(shp2)
                    if sk2:
                        base2["shape_kind"] = sk2

                    # Minimal text capture for placeholders/layout labels
                    if kind2 == "text" and getattr(shp2, "has_text_frame", False) and shp2.has_text_frame:
                        try:
                            t2 = _norm_text(shp2.text_frame.text or "")
                        except Exception:
                            t2 = ""
                        if t2:
                            base2["text"] = t2
                            base2["normalized_text"] = t2
                            base2["text_struct"] = _text_structure(shp2, theme_rgb)

                    elif kind2 == "image":
                        base2["text"] = "[IMAGE]"
                        base2["normalized_text"] = "[IMAGE]"
                        base2["image"] = _image_info(shp2)

                    elif kind2 == "table":
                        base2["text"] = "[TABLE]"
                        base2["normalized_text"] = "[TABLE]"
                        base2["table"] = _table_info(shp2, theme_rgb)

                    chunks.append(base2)
                    z += 1
                return z

            z0 = 1
            try:
                z0 = _emit_from_shape_collection(slide.slide_layout.shapes, "ly", z0)
            except Exception:
                pass
            try:
                _emit_from_shape_collection(slide.slide_master.shapes, "ms", z0)
            except Exception:
                pass

        for shape_idx, shp in enumerate(slide.shapes, start=1):
            kind = _shape_kind(shp)

            bbox = {
                "x": _emu(getattr(shp, "left", 0)),
                "y": _emu(getattr(shp, "top", 0)),
                "w": _emu(getattr(shp, "width", 0)),
                "h": _emu(getattr(shp, "height", 0)),
            }

            # Rotation is degrees in python-pptx (may be None)
            try:
                rot = float(getattr(shp, "rotation", 0) or 0)
            except Exception:
                rot = 0.0

            base: Dict[str, Any] = {
                "chunk_id": f"s{slide_idx:03d}_sh{shape_idx:03d}",
                "page": slide_idx,
                "bbox": bbox,
                "text": "",
                "normalized_text": "",
                "heading_level": 0,
            }

            # Optional extended fields (only when schema/run expects them)
            if include_extended:
                base["kind"] = kind
                base["z"] = shape_idx + 10000
                base["rotation_deg"] = rot
                base["style"] = {
                    "fill": _fill_dict(shp, theme_rgb),
                    "line": _line_dict(shp, theme_rgb),
                }
                sk = _autoshape_kind(shp)
                if sk:
                    base["shape_kind"] = sk

            # Text
            if kind == "text":
                try:
                    text = _norm_text(shp.text_frame.text or "")
                except Exception:
                    text = ""
                if text:
                    base["text"] = text
                    base["normalized_text"] = text
                if include_extended:
                    base["text_struct"] = _text_structure(shp, theme_rgb)

            # Table
            elif kind == "table":
                if include_extended:
                    base["table"] = _table_info(shp, theme_rgb)
                # Keep something in text fields for downstream that expects non-empty sometimes
                # (schema itself may not require it; safe default)
                base["text"] = "[TABLE]"
                base["normalized_text"] = "[TABLE]"

            # Image
            elif kind == "image":
                if include_extended:
                    base["image"] = _image_info(shp)
                base["text"] = "[IMAGE]"
                base["normalized_text"] = "[IMAGE]"

            # Group / Shape / Other
            else:
                # If it has text frame even though classified otherwise, capture it.
                if getattr(shp, "has_text_frame", False) and shp.has_text_frame:
                    try:
                        text = _norm_text(shp.text_frame.text or "")
                    except Exception:
                        text = ""
                    if text:
                        base["text"] = text
                        base["normalized_text"] = text
                        if include_extended:
                            base["text_struct"] = _text_structure(shp, theme_rgb)

            # Keep all shapes; downstream can filter by kind
            chunks.append(base)

    document: Dict[str, Any] = {
        "document_id": document_id,
        "source_type": "pptx",
        "source_path": str(p),
        "page_count": len(prs.slides),
    }
    if include_extended:
        document["page"] = page

    return {
        "schema_version": "0.1",
        "document": document,
        "chunks": chunks,
    }

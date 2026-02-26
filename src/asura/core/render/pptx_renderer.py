from __future__ import annotations

import json
import copy
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

# Helper functions for no-fill and no-line
def _set_no_fill(shape: Any) -> None:
    try:
        spPr = shape._element.spPr
    except Exception:
        return
    try:
        for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill", "a:grpFill"):
            el = spPr.find(qn(tag))
            if el is not None:
                spPr.remove(el)
    except Exception:
        pass
    try:
        spPr.append(OxmlElement("a:noFill"))
    except Exception:
        pass


def _set_no_line(shape: Any) -> None:
    # 1) Best-effort: use python-pptx API to neutralize theme-default outlines.
    try:
        ln_api = getattr(shape, "line", None)
        if ln_api is not None:
            try:
                ln_api.width = Pt(0)
            except Exception:
                pass
            try:
                # Prefer background fill for the line when available.
                ln_api.fill.background()
            except Exception:
                pass
            try:
                # Fully transparent as an extra guard.
                ln_api.color.transparency = 1.0
            except Exception:
                pass
    except Exception:
        pass

    # 2) Authoritative: enforce DrawingML <a:ln w="0"><a:noFill/></a:ln>
    try:
        spPr = shape._element.spPr
    except Exception:
        return

    # 2a) Kill theme-default outline refs in p:style/a:lnRef by setting idx="0".
    # NOTE: p:style requires all 4 children (lnRef, fillRef, effectRef, fontRef).
    # Removing lnRef would violate the OOXML schema and trigger PowerPoint's repair prompt.
    # Setting idx="0" is the correct OOXML way to indicate "no line style from theme matrix".
    try:
        el = shape._element
        st_el = el.find(qn("p:style"))
        if st_el is not None:
            ln_ref = st_el.find(qn("a:lnRef"))
            if ln_ref is not None:
                ln_ref.set("idx", "0")
    except Exception:
        pass

    ln = None
    try:
        ln = spPr.find(qn("a:ln"))
    except Exception:
        ln = None
    if ln is None:
        try:
            ln = OxmlElement("a:ln")
            spPr.append(ln)
        except Exception:
            return
    try:
        ln.set("w", "0")
    except Exception:
        pass
    try:
        for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:pattFill"):
            el = ln.find(qn(tag))
            if el is not None:
                ln.remove(el)
    except Exception:
        pass
    try:
        ln.append(OxmlElement("a:noFill"))
    except Exception:
        pass

    # Extra: remove dash/join/cap that could survive in some viewers.
    try:
        for tag in ("a:prstDash", "a:round", "a:bevel", "a:miter", "a:headEnd", "a:tailEnd"):
            el2 = ln.find(qn(tag))
            if el2 is not None:
                ln.remove(el2)
    except Exception:
        pass


# XML-based solid fill alpha setter
def _set_solid_fill_alpha_xml(shape: Any, alpha01: float) -> None:
    """Set alpha (opacity) for a solid fill using DrawingML XML.

    alpha01: 0..1 where 1=fully opaque.
    """
    try:
        a = float(alpha01)
    except Exception:
        return
    if a < 0.0:
        a = 0.0
    if a > 1.0:
        a = 1.0

    try:
        spPr = shape._element.spPr
    except Exception:
        return

    # Find a:solidFill
    try:
        solid = spPr.find(qn("a:solidFill"))
    except Exception:
        solid = None
    if solid is None:
        return

    # Find a:srgbClr inside solidFill
    try:
        clr = solid.find(qn("a:srgbClr"))
    except Exception:
        clr = None
    if clr is None:
        return

    # Remove existing alpha modifiers
    try:
        for tag in ("a:alpha", "a:alphaMod", "a:alphaOff"):
            el = clr.find(qn(tag))
            if el is not None:
                clr.remove(el)
    except Exception:
        pass

    # a:alpha val is 0..100000 (percent * 1000)
    try:
        ael = OxmlElement("a:alpha")
        ael.set("val", str(int(round(a * 100000))))
        clr.append(ael)
    except Exception:
        pass

import base64
from io import BytesIO
import math
from pptx.parts.image import Image as PptxImage
import hashlib
import zipfile



PT_PER_INCH = 72.0


EMU_PER_INCH = 914400.0


def _emu_to_pt(x_emu: Any) -> float:
    try:
        return float(x_emu) * PT_PER_INCH / EMU_PER_INCH
    except Exception:
        return 0.0


def _bbox_emu_to_pt_rect(bbox: Any) -> tuple[float, float, float, float] | None:
    """Convert bbox in EMU dict {x,y,w,h} to (x1,y1,x2,y2) in points."""
    if not isinstance(bbox, dict):
        return None
    if not all(k in bbox for k in ("x", "y", "w", "h")):
        return None
    x = _emu_to_pt(bbox.get("x"))
    y = _emu_to_pt(bbox.get("y"))
    w = _emu_to_pt(bbox.get("w"))
    h = _emu_to_pt(bbox.get("h"))
    return (x, y, x + w, y + h)


def _align_from_any(v: Any) -> PP_ALIGN | None:
    if v is None:
        return None
    # python-pptx often stringifies enums like "CENTER (2)" or "PP_ALIGN.CENTER (2)"
    if isinstance(v, PP_ALIGN):
        return v
    if isinstance(v, (int, float)):
        try:
            return PP_ALIGN(int(v))
        except Exception:
            return None
    if not isinstance(v, str):
        return None

    s = v.strip().lower()
    if not s:
        return None

    # Exact-ish tokens
    if s in ("left", "l"):
        return PP_ALIGN.LEFT
    if s in ("center", "centre", "c"):
        return PP_ALIGN.CENTER
    if s in ("right", "r"):
        return PP_ALIGN.RIGHT
    if s in ("justify", "justified", "j"):
        return PP_ALIGN.JUSTIFY

    # Substring heuristics for values like "center (2)", "pp_align.center (2)", etc.
    if "center" in s or "centre" in s:
        return PP_ALIGN.CENTER
    if "right" in s:
        return PP_ALIGN.RIGHT
    if "left" in s:
        return PP_ALIGN.LEFT
    if "justify" in s:
        return PP_ALIGN.JUSTIFY

    return None


def _alpha01(v: Any) -> float | None:
    """Normalize alpha into 0..1.

    Accepts:
    - 0..1 floats
    - 0..100 percentages
    - 0..255 byte alpha
    Returns None if v is not numeric.
    """
    if isinstance(v, str):
        s = v.strip()
        if not s:
            return None
        try:
            v = float(s)
        except Exception:
            return None
    if not isinstance(v, (int, float)):
        return None
    try:
        x = float(v)
    except Exception:
        return None
    if x < 0:
        x = 0.0
    # Heuristics for common encodings
    if x > 1.0:
        if x <= 100.0:
            x = x / 100.0
        elif x <= 255.0:
            x = x / 255.0
        else:
            x = 1.0
    if x < 0.0:
        x = 0.0
    if x > 1.0:
        x = 1.0
    return x


def _vanchor_from_any(v: Any) -> MSO_VERTICAL_ANCHOR | None:
    if v is None:
        return None
    if isinstance(v, MSO_VERTICAL_ANCHOR):
        return v
    if isinstance(v, (int, float)):
        try:
            return MSO_VERTICAL_ANCHOR(int(v))
        except Exception:
            return None
    if not isinstance(v, str):
        return None
    s = v.strip().lower()
    if not s:
        return None
    if s in ("top", "t"):
        return MSO_VERTICAL_ANCHOR.TOP
    if s in ("middle", "center", "centre", "mid", "m"):
        return MSO_VERTICAL_ANCHOR.MIDDLE
    if s in ("bottom", "b"):
        return MSO_VERTICAL_ANCHOR.BOTTOM
    if "middle" in s or "center" in s or "centre" in s:
        return MSO_VERTICAL_ANCHOR.MIDDLE
    if "bottom" in s:
        return MSO_VERTICAL_ANCHOR.BOTTOM
    if "top" in s:
        return MSO_VERTICAL_ANCHOR.TOP
    return None


def _style_from_extraction(style: Any) -> dict[str, Any]:
    """Best-effort adapter from extraction.extended style to renderer's minimal style keys."""
    out: dict[str, Any] = {}
    if not isinstance(style, dict):
        return out

    # Fill: expect style.fill = {type: 'solid'|'none', color_rgb: 'RRGGBB' or '#RRGGBB'}
    fill = style.get("fill")
    if isinstance(fill, dict):
        ftype = str(fill.get("type", "")).lower()
        if ftype in ("none", "transparent"):
            out["fill_none"] = True
        else:
            c = fill.get("color_rgb")
            if isinstance(c, str) and c:
                out["fill_rgb"] = c
            a = fill.get("alpha")
            if a is None:
                a = fill.get("opacity")
            if a is None:
                a = fill.get("opacity_pct")
            tr = fill.get("transparency")
            if a is None and tr is not None:
                tt = _alpha01(tr)
                if tt is not None:
                    a = 1.0 - tt
            aa = _alpha01(a)
            if aa is not None:
                out["fill_alpha"] = aa

    # Line: expect style.line = {type, visible, color_rgb, width_emu, alpha}
    line = style.get("line")
    if isinstance(line, dict):
        ltype = str(line.get("type") or "").lower().strip()
        vis = line.get("visible")
        w_emu = line.get("width_emu")
        a = line.get("alpha")
        if a is None:
            a = line.get("opacity")
        if a is None:
            a = line.get("opacity_pct")
        tr = line.get("transparency")
        if a is None and tr is not None:
            tt = _alpha01(tr)
            if tt is not None:
                a = 1.0 - tt

        # If extractor explicitly marks no-line, honor it.
        if ltype in ("none", "no", "transparent") or vis is False:
            out["line_none"] = True
        else:
            # Width <= 0 means no visible line.
            try:
                if isinstance(w_emu, (int, float)) and float(w_emu) <= 0:
                    out["line_none"] = True
            except Exception:
                pass
            # Alpha <= 0 means fully transparent line.
            try:
                if isinstance(a, (int, float)) and float(a) <= 0:
                    out["line_none"] = True
            except Exception:
                pass

        # Only emit line style when line is not none.
        if out.get("line_none") is not True:
            c = line.get("color_rgb")
            if isinstance(c, str) and c:
                out["line_rgb"] = c
            if isinstance(w_emu, (int, float)):
                try:
                    w_pt = _emu_to_pt(w_emu)
                    if w_pt > 0:
                        out["line_width_pt"] = w_pt
                except Exception:
                    pass
            aa = _alpha01(a)
            if aa is not None:
                out["line_alpha"] = aa

    return out


def _rgb_from_any(v: Any) -> RGBColor | None:
    """Parse RGB from '#RRGGBB' or 'RRGGBB' or [r,g,b]."""
    if v is None:
        return None
    if isinstance(v, (list, tuple)) and len(v) == 3:
        try:
            r, g, b = (int(v[0]), int(v[1]), int(v[2]))
            if 0 <= r <= 255 and 0 <= g <= 255 and 0 <= b <= 255:
                return RGBColor(r, g, b)
        except Exception:
            return None
        return None
    if isinstance(v, str):
        s = v.strip()
        if s.startswith("#"):
            s = s[1:]
        if len(s) == 6:
            try:
                r = int(s[0:2], 16)
                g = int(s[2:4], 16)
                b = int(s[4:6], 16)
                return RGBColor(r, g, b)
            except Exception:
                return None
    return None


# Helper: build sha256->(bytes, ext) map from pptx media
def _build_pptx_media_sha_map(pptx_path: Path) -> dict[str, tuple[bytes, str]]:
    """Scan ppt/media/* inside a .pptx and map sha256 -> (bytes, ext).

    This allows DOM rendering to place images even when extractor did not materialize assets
    into run_dir/assets.
    """
    out: dict[str, tuple[bytes, str]] = {}
    try:
        with zipfile.ZipFile(pptx_path, "r") as zf:
            for name in zf.namelist():
                if not name.startswith("ppt/media/"):
                    continue
                try:
                    blob = zf.read(name)
                except Exception:
                    continue
                if not blob:
                    continue
                h = hashlib.sha256(blob).hexdigest()
                ext = Path(name).suffix.lower().lstrip(".") or "png"
                out[h] = (blob, ext)
    except Exception:
        return {}
    return out



def _find_source_pptx(run_dir: Path) -> Path | None:
    """Best-effort discovery of the original .pptx used to produce extraction.json."""

    # 1) runlog.json (top-level or inputs)
    try:
        runlog_path = run_dir / "runlog.json"
        if runlog_path.exists():
            runlog = _load_json(runlog_path)
            if isinstance(runlog, dict):
                candidates: list[str] = []
                for k in ("input_path", "source_path", "input", "source"):
                    v = runlog.get(k)
                    if isinstance(v, str) and v:
                        candidates.append(v)
                inputs = runlog.get("inputs")
                if isinstance(inputs, dict):
                    for k in ("input_path", "source_path", "input", "source"):
                        v = inputs.get(k)
                        if isinstance(v, str) and v:
                            candidates.append(v)

                for src in candidates:
                    if not src.lower().endswith(".pptx"):
                        continue
                    p = Path(src)
                    if not p.is_absolute():
                        # try relative to run_dir and project root-ish
                        for base in (run_dir, run_dir.parent, run_dir.parent.parent):
                            pp = (base / p).resolve()
                            if pp.exists():
                                return pp
                    else:
                        if p.exists():
                            return p
    except Exception:
        pass

    # 2) common project locations
    common = [
        run_dir / "source.pptx",
        run_dir / "assets" / "source.pptx",
        run_dir.parent / "input" / "テスト2.pptx",
        run_dir.parent.parent / "input" / "テスト2.pptx",
        run_dir.parent.parent.parent / "input" / "テスト2.pptx",
    ]
    for p in common:
        try:
            if p.exists() and p.suffix.lower() == ".pptx":
                return p.resolve()
        except Exception:
            continue

    # 3) heuristic scan upward for an `input` directory
    try:
        cur = run_dir
        for _ in range(8):
            input_dir = cur / "input"
            if input_dir.exists() and input_dir.is_dir():
                pptxs = sorted(input_dir.glob("*.pptx"), key=lambda x: x.stat().st_mtime, reverse=True)
                if pptxs:
                    return pptxs[0].resolve()
            cur = cur.parent
    except Exception:
        pass

    return None


def _parse_chunk_id_slide_shape(chunk_id: Any) -> tuple[int, int] | None:
    """Parse chunk_id like 's006_sh009' into (slide_no=6, shape_ord=9)."""
    if not isinstance(chunk_id, str):
        return None
    s = chunk_id.strip()
    if not s.startswith("s"):
        return None
    try:
        parts = s.split("_")
        if len(parts) < 2:
            return None
        slide_no = int(parts[0][1:])
        sh_part = None
        for p in parts[1:]:
            if p.startswith("sh") and len(p) >= 4:
                sh_part = p
                break
        if sh_part is None:
            return None
        shape_no = int(sh_part[2:])
        if slide_no <= 0 or shape_no <= 0:
            return None
        return (slide_no, shape_no)
    except Exception:
        return None


# Copy avLst (adjustment values) from source shape to destination shape.
def _copy_adj_from_src(*, prs_src: Any, slide_no: int, shape_no: int, dst_shape: Any) -> bool:
    """Copy <a:avLst> contents from source shape's prstGeom into destination shape.

    Preserves corner radius (adj) for roundRect and other parameterized presets.
    Returns True if any gd elements were copied.
    """
    try:
        src_slide = prs_src.slides[slide_no - 1]
        src_shape = src_slide.shapes[shape_no - 1]
    except Exception:
        return False

    try:
        src_spPr = src_shape._element.spPr
        src_geom = src_spPr.find(qn("a:prstGeom"))
        if src_geom is None:
            return False
        src_avLst = src_geom.find(qn("a:avLst"))
        if src_avLst is None or len(list(src_avLst)) == 0:
            return False
    except Exception:
        return False

    try:
        dst_spPr = dst_shape._element.spPr
        dst_geom = dst_spPr.find(qn("a:prstGeom"))
        if dst_geom is None:
            return False
        dst_avLst = dst_geom.find(qn("a:avLst"))
        if dst_avLst is None:
            return False
    except Exception:
        return False

    try:
        for child in list(dst_avLst):
            dst_avLst.remove(child)
        for gd in src_avLst:
            dst_avLst.append(parse_xml(gd.xml))
        return True
    except Exception:
        return False


# Copy line XML from source shape to destination shape.
def _copy_line_xml_from_src(*, prs_src: Any, slide_no: int, shape_no: int, dst_shape: Any) -> bool:
    """Copy <a:ln> from source shape into destination shape.

    This preserves theme-based outlines (schemeClr, presets, dash, etc.) that python-pptx may not expose.
    Returns True if a:ln was copied.
    """
    try:
        src_slide = prs_src.slides[slide_no - 1]
        src_shape = src_slide.shapes[shape_no - 1]
    except Exception:
        return False

    try:
        src_spPr = src_shape._element.spPr
        src_ln = src_spPr.find(qn("a:ln"))
        if src_ln is None:
            return False
    except Exception:
        return False

    try:
        dst_spPr = dst_shape._element.spPr
    except Exception:
        return False

    # Remove existing a:ln in destination
    try:
        dst_ln = dst_spPr.find(qn("a:ln"))
        if dst_ln is not None:
            dst_spPr.remove(dst_ln)
    except Exception:
        pass

    # Clone via parse_xml to keep namespaces/structure valid for python-pptx.
    try:
        dst_spPr.append(parse_xml(src_ln.xml))
        return True
    except Exception:
        return False


def _resolve_textbox_fill_from_src(prs_src: Any, slide_no: int, shape_no: int) -> dict[str, Any]:
    """Best-effort resolve textbox fill from the source pptx (theme-resolved when possible).

    Returns style keys compatible with _style_from_extraction output: {fill_none, fill_rgb, fill_alpha}.
    """
    out: dict[str, Any] = {}
    try:
        slide = prs_src.slides[slide_no - 1]
        shp = slide.shapes[shape_no - 1]
    except Exception:
        return out

    # Only attempt when it looks like a textbox/text frame shape.
    try:
        if not getattr(shp, "has_text_frame", False):
            return out
    except Exception:
        return out

    try:
        fill = shp.fill
        ftype = fill.type
    except Exception:
        return out

    try:
        if ftype in (MSO_FILL.BACKGROUND,):
            # BACKGROUND means same as slide background; visually, this should be transparent.
            out["fill_none"] = True
            return out

        if ftype == MSO_FILL.SOLID:
            try:
                fill_rgb = _rgb_from_any(getattr(fill.fore_color, "rgb", None))
            except Exception:
                fill_rgb = None

            # Some theme cases still expose .rgb; if not, we cannot do much here.
            if fill_rgb is not None:
                # Convert RGBColor -> 'RRGGBB' string for our style dict
                try:
                    out["fill_rgb"] = "%02X%02X%02X" % (fill_rgb[0], fill_rgb[1], fill_rgb[2])
                except Exception:
                    pass

            # transparency: 0..1 (1=fully transparent)
            try:
                tr = getattr(fill.fore_color, "transparency", None)
                if tr is not None:
                    a = 1.0 - float(tr)
                    if a < 0:
                        a = 0.0
                    if a > 1:
                        a = 1.0
                    out["fill_alpha"] = a
            except Exception:
                pass

            return out

        # For other fill types (gradient/picture/pattern), we currently do not resolve.
        return out
    except Exception:
        return out


def _resolve_textbox_alignments_from_src(prs_src: Any, slide_no: int, shape_no: int) -> list[PP_ALIGN | None]:
    """Resolve paragraph alignments from the source PPTX textbox.

    Returns a list aligned to prs_src.text_frame.paragraphs order.
    """
    out: list[PP_ALIGN | None] = []
    try:
        slide = prs_src.slides[slide_no - 1]
        shp = slide.shapes[shape_no - 1]
    except Exception:
        return out

    try:
        if not getattr(shp, "has_text_frame", False):
            return out
    except Exception:
        return out

    try:
        tf = shp.text_frame
        for p in tf.paragraphs:
            try:
                out.append(p.alignment)
            except Exception:
                out.append(None)
    except Exception:
        return []

    return out


def _resolve_textbox_vanchor_from_src(prs_src: Any, slide_no: int, shape_no: int) -> MSO_VERTICAL_ANCHOR | None:
    """Resolve text frame vertical anchor (top/middle/bottom) from source PPTX."""
    try:
        slide = prs_src.slides[slide_no - 1]
        shp = slide.shapes[shape_no - 1]
    except Exception:
        return None

    try:
        if not getattr(shp, "has_text_frame", False):
            return None
    except Exception:
        return None

    try:
        return shp.text_frame.vertical_anchor
    except Exception:
        return None


# --- New: Shape style resolver for fill/line alpha from source PPTX ---
def _resolve_shape_style_from_src(prs_src: Any, slide_no: int, shape_no: int) -> dict[str, Any]:
    """Resolve basic shape style (fill/line rgb + alpha + none flags) from the source PPTX.

    This is used when extraction lacks transparency information.
    Returns keys compatible with `_style_from_extraction` output.
    """
    out: dict[str, Any] = {}
    try:
        slide = prs_src.slides[slide_no - 1]
        shp = slide.shapes[shape_no - 1]
    except Exception:
        return out

    # --- Fill ---
    try:
        fill = shp.fill
        ftype = fill.type
        if ftype in (MSO_FILL.BACKGROUND,):
            out["fill_none"] = True
        elif ftype == MSO_FILL.SOLID:
            try:
                rgb = getattr(fill.fore_color, "rgb", None)
                rr = _rgb_from_any(rgb)
                if rr is not None:
                    out["fill_rgb"] = "%02X%02X%02X" % (rr[0], rr[1], rr[2])
            except Exception:
                pass
            # python-pptx transparency: 0..1 (1=fully transparent)
            try:
                tr = getattr(fill.fore_color, "transparency", None)
                if tr is not None:
                    a = 1.0 - float(tr)
                    if a < 0.0:
                        a = 0.0
                    if a > 1.0:
                        a = 1.0
                    out["fill_alpha"] = a
            except Exception:
                pass
        else:
            # gradients/pictures/patterns: leave as-is
            pass
    except Exception:
        pass

    # --- Line ---
    try:
        ln = shp.line
        # If width is 0 or None, treat as no line.
        try:
            w = getattr(ln, "width", None)
            if w is not None and float(w.pt) <= 0:
                out["line_none"] = True
        except Exception:
            pass

        # Color
        try:
            rgb = getattr(getattr(ln, "color", None), "rgb", None)
            rr = _rgb_from_any(rgb)
            if rr is not None:
                out["line_rgb"] = "%02X%02X%02X" % (rr[0], rr[1], rr[2])
        except Exception:
            pass

        # Transparency
        try:
            tr = getattr(getattr(ln, "color", None), "transparency", None)
            if tr is not None:
                a = 1.0 - float(tr)
                if a < 0.0:
                    a = 0.0
                if a > 1.0:
                    a = 1.0
                out["line_alpha"] = a
        except Exception:
            pass

        # Width (pt)
        try:
            w = getattr(ln, "width", None)
            if w is not None:
                wpt = float(w.pt)
                if wpt > 0:
                    out["line_width_pt"] = wpt
        except Exception:
            pass

        # If line is effectively invisible (alpha == 0), mark none.
        try:
            if isinstance(out.get("line_alpha"), (int, float)) and float(out.get("line_alpha")) <= 0:
                out["line_none"] = True
        except Exception:
            pass

        # XML fallback: detect presence/width/noFill even when theme-based color is not exposed.
        try:
            spPr = shp._element.spPr
            xln = spPr.find(qn("a:ln"))
            if xln is not None:
                # noFill means no visible line.
                if xln.find(qn("a:noFill")) is not None:
                    out["line_none"] = True
                # solidFill with alpha=0 is also fully transparent (invisible).
                # python-pptx's color.transparency often fails for line colors, so check XML directly.
                if not out.get("line_none"):
                    solid = xln.find(qn("a:solidFill"))
                    if solid is not None:
                        for ctag in (qn("a:srgbClr"), qn("a:schemeClr"), qn("a:prstClr")):
                            clr = solid.find(ctag)
                            if clr is not None:
                                ael = clr.find(qn("a:alpha"))
                                if ael is not None:
                                    try:
                                        if int(ael.get("val", "100000")) <= 0:
                                            out["line_none"] = True
                                            out.pop("line_rgb", None)
                                            out.pop("line_width_pt", None)
                                    except Exception:
                                        pass
                                break
                # Width attribute is in EMU (1pt = 12700 EMU). Only set if line is visible.
                if not out.get("line_none"):
                    w_attr = xln.get("w")
                    if w_attr is not None:
                        try:
                            w_emu = float(w_attr)
                            w_pt = _emu_to_pt(w_emu)
                            if w_pt > 0:
                                out["line_width_pt"] = w_pt
                        except Exception:
                            pass
        except Exception:
            pass
    except Exception:
        pass

    return out


def _apply_common_style(shape: Any, el: dict[str, Any]) -> None:
    """Apply minimal style fields if present.

    Supported keys (all optional):
    - rotation_deg: float
    - style.fill_rgb: '#RRGGBB' | 'RRGGBB' | [r,g,b]
    - style.fill_none: bool
    - style.line_rgb: '#RRGGBB' | 'RRGGBB' | [r,g,b]
    - style.line_width_pt: float
    """
    rot = el.get("rotation_deg")
    if rot is not None:
        try:
            shape.rotation = float(rot)
        except Exception:
            pass

    style = el.get("style")
    if not isinstance(style, dict):
        return

    # Suppress default theme effects (e.g. drop shadow from effectRef idx="2" in autoshape template).
    try:
        _st_el = shape._element.find(qn("p:style"))
        if _st_el is not None:
            _eff_ref = _st_el.find(qn("a:effectRef"))
            if _eff_ref is not None:
                _eff_ref.set("idx", "0")
    except Exception:
        pass
    # Also strip any explicit effectLst from spPr (defensive; renderer never adds effects).
    try:
        _spPr = shape._element.spPr
        _eff_lst = _spPr.find(qn("a:effectLst"))
        if _eff_lst is not None:
            _spPr.remove(_eff_lst)
    except Exception:
        pass

    # Fill
    try:
        if style.get("fill_none") is True:
            _set_no_fill(shape)
        else:
            fill_rgb = _rgb_from_any(style.get("fill_rgb"))
            if fill_rgb is not None:
                shape.fill.solid()
                shape.fill.fore_color.rgb = fill_rgb
                a = style.get("fill_alpha")
                if a is not None:
                    aa = _alpha01(a)
                    if aa is not None:
                        # Prefer XML alpha because python-pptx transparency setters can be unreliable.
                        try:
                            _set_solid_fill_alpha_xml(shape, aa)
                        except Exception:
                            pass
                        # Best-effort also set python-pptx transparency if available.
                        try:
                            shape.fill.fore_color.transparency = 1.0 - aa
                        except Exception:
                            pass
    except Exception:
        pass

    try:
        # Only draw a line when we have an explicit line spec.
        # Otherwise python-pptx may leave the default theme outline (often blue).
        has_explicit_line = (
            ("line_rgb" in style)
            or ("line_width_pt" in style)
            or ("line_alpha" in style)
            or (style.get("line_from_src") is True)
        )

        if style.get("line_none") is True:
            _set_no_line(shape)
            return

        if not has_explicit_line:
            _set_no_line(shape)
            return

        if style.get("line_from_src") is True and not any(
            k in style for k in ("line_rgb", "line_width_pt", "line_alpha")
        ):
            # Keep whatever is already on the shape (typically copied from the source XML).
            return

        # Reset lnRef to prevent theme's line color (typically accent1 = blue) from leaking
        # when explicit line properties are set without an explicit color.
        try:
            _st_el = shape._element.find(qn("p:style"))
            if _st_el is not None:
                _ln_ref = _st_el.find(qn("a:lnRef"))
                if _ln_ref is not None:
                    _ln_ref.set("idx", "0")
        except Exception:
            pass

        line_rgb = _rgb_from_any(style.get("line_rgb"))
        if line_rgb is not None:
            shape.line.color.rgb = line_rgb

        a = style.get("line_alpha")
        if a is not None:
            aa = _alpha01(a)
            if aa is not None:
                try:
                    shape.line.color.transparency = 1.0 - aa
                except Exception:
                    pass

        lw = style.get("line_width_pt")
        if lw is not None:
            shape.line.width = Pt(float(lw))
        else:
            # If width is missing but line is explicit via color, keep a sane minimal width.
            try:
                shape.line.width = Pt(0.75)
            except Exception:
                pass
    except Exception:
        pass


def _render_table_cells(table: Any, cells: list[Any]) -> None:
    """Render table cells with per-run font info when available, plain text otherwise.

    cell_spec keys: r, c, text, paragraphs (optional list of paragraph dicts with runs).
    """
    for cell_spec in cells:
        if not isinstance(cell_spec, dict):
            continue
        r = cell_spec.get("r")
        c = cell_spec.get("c")
        text = cell_spec.get("text")
        if not isinstance(r, int) or not isinstance(c, int):
            continue
        try:
            cell_obj = table.cell(r, c)
        except Exception:
            continue

        paragraphs_data = cell_spec.get("paragraphs")
        if isinstance(paragraphs_data, list) and paragraphs_data:
            try:
                tf = cell_obj.text_frame
                tf.clear()
                first_para = True
                for p_spec in paragraphs_data:
                    if not isinstance(p_spec, dict):
                        continue
                    runs = p_spec.get("runs") or []
                    para = tf.paragraphs[0] if first_para else tf.add_paragraph()
                    first_para = False
                    al = _align_from_any(
                        p_spec.get("alignment") or p_spec.get("align") or p_spec.get("text_align")
                    )
                    if al is not None:
                        try:
                            para.alignment = al
                        except Exception:
                            pass
                    for r_spec in runs:
                        if not isinstance(r_spec, dict):
                            continue
                        txt = r_spec.get("text")
                        if txt is None:
                            continue
                        run = para.add_run()
                        run.text = str(txt)
                        try:
                            fn = r_spec.get("font_name")
                            if isinstance(fn, str) and fn:
                                run.font.name = fn
                        except Exception:
                            pass
                        try:
                            fs = r_spec.get("font_size_emu")
                            if isinstance(fs, (int, float)) and fs > 0:
                                run.font.size = Pt(_emu_to_pt(fs))
                        except Exception:
                            pass
                        try:
                            b = r_spec.get("bold")
                            if b is True:
                                run.font.bold = True
                            elif b is False:
                                run.font.bold = False
                        except Exception:
                            pass
                        try:
                            if r_spec.get("italic") is True:
                                run.font.italic = True
                        except Exception:
                            pass
                        try:
                            if r_spec.get("underline") is True:
                                run.font.underline = True
                        except Exception:
                            pass
                        try:
                            rgb = _rgb_from_any(r_spec.get("color_rgb"))
                            if rgb is not None:
                                run.font.color.rgb = rgb
                        except Exception:
                            pass
            except Exception:
                # Fallback to plain text if paragraph rendering fails.
                try:
                    cell_obj.text = "" if text is None else str(text)
                except Exception:
                    pass
        else:
            try:
                cell_obj.text = "" if text is None else str(text)
            except Exception:
                pass


def _autoshape_type_from_el(el: dict[str, Any]) -> MSO_AUTO_SHAPE_TYPE:
    """Map our schema-ish shape type into a pptx autoshape."""
    st_raw = str(el.get("shape_type", "rect"))
    st = st_raw.strip().lower()
    st_norm = st.replace("-", "_").replace(" ", "_")
    if st_norm in (
        "round_rect",
        "roundrect",
        "round_rectangle",
        "rounded_rect",
        "roundedrect",
        "rounded_rectangle",
        "roundedrectangle",
        "roundrectangle",
        "roundrectangular",
        "roundedrectangular",
    ):
        return MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
    if st_norm in ("ellipse", "oval", "circle"):
        return MSO_AUTO_SHAPE_TYPE.OVAL
    # default
    return MSO_AUTO_SHAPE_TYPE.RECTANGLE


def _pt_to_inches(x_pt: float) -> Inches:
    return Inches(x_pt / PT_PER_INCH)


# Coerce to float points
def _pt(x: Any) -> float:
    try:
        return float(x)
    except Exception:
        return 0.0
def _add_picture_with_fit(
    *,
    slide: Any,
    run_dir: Path,
    left_pt: float,
    top_pt: float,
    box_w_pt: float,
    box_h_pt: float,
    img_spec: Any,
    fit: str = "stretch",
) -> Any | None:
    """Add an image into a bbox with fit policy.

    - stretch: force to bbox (may distort)
    - contain: preserve aspect, fit inside bbox (letterbox)
    - cover: preserve aspect, fill bbox (crop)

    img_spec can be:
      - {"path": "..."} (relative paths resolved from run_dir)
      - {"bytes_b64": "..."}
      - "path/to/file"
    """

    # Resolve image bytes or path
    img_path: Path | None = None
    blob: bytes | None = None

    if isinstance(img_spec, dict):
        # Common path keys emitted by extractors
        p = (
            img_spec.get("path")
            or img_spec.get("relpath")
            or img_spec.get("relative_path")
            or img_spec.get("file")
            or img_spec.get("filename")
            or img_spec.get("name")
        )

        # Common base64 keys
        b64 = (
            img_spec.get("bytes_b64")
            or img_spec.get("data_b64")
            or img_spec.get("blob_b64")
        )

        if isinstance(p, str) and p:
            img_path = Path(p)
            if not img_path.is_absolute():
                img_path = (run_dir / img_path).resolve()

        if blob is None and isinstance(b64, str) and b64:
            try:
                blob = base64.b64decode(b64)
            except Exception:
                blob = None

        # If extractor provides content-addressed asset info, try common locations.
        if img_path is None:
            sha = img_spec.get("sha256")
            ext = img_spec.get("ext") or img_spec.get("format")
            if isinstance(sha, str) and sha:
                # normalize extension like '.png' vs 'png'
                if isinstance(ext, str) and ext:
                    e = ext.strip().lower()
                    if e.startswith("."):
                        e = e[1:]
                else:
                    e = "png"
                candidates = [
                    run_dir / "assets" / f"{sha}.{e}",
                    run_dir / f"{sha}.{e}",
                    run_dir / "images" / f"{sha}.{e}",
                ]
                for cand in candidates:
                    if cand.exists():
                        img_path = cand.resolve()
                        break
    elif isinstance(img_spec, str) and img_spec:
        img_path = Path(img_spec)
        if not img_path.is_absolute():
            img_path = (run_dir / img_path).resolve()
    elif isinstance(img_spec, (bytes, bytearray)) and img_spec:
        blob = bytes(img_spec)

    if img_path is not None and not img_path.exists():
        img_path = None

    if img_path is None and blob is None:
        return None

    # Get intrinsic pixel size
    try:
        if blob is not None:
            im = PptxImage.from_blob(blob)
        else:
            im = PptxImage.from_file(str(img_path))  # type: ignore[arg-type]
        iw = float(getattr(im, "px_width"))
        ih = float(getattr(im, "px_height"))
        if iw <= 0 or ih <= 0:
            raise ValueError("invalid image px size")
    except Exception:
        # Fallback: if we cannot read size, just stretch.
        iw, ih = 1.0, 1.0
        fit = "stretch"

    fit_n = str(fit or "stretch").lower().strip()
    if fit_n not in ("stretch", "contain", "cover"):
        fit_n = "stretch"

    # Build base args
    left = Pt(left_pt)
    top = Pt(top_pt)
    box_w = Pt(box_w_pt)
    box_h = Pt(box_h_pt)

    # Helper to call add_picture with either path or blob
    def _add_picture(*, l: Pt, t: Pt, w: Pt | None, h: Pt | None) -> Any | None:
        try:
            if blob is not None:
                stream = BytesIO(blob)
                return slide.shapes.add_picture(stream, l, t, width=w, height=h)
            return slide.shapes.add_picture(str(img_path), l, t, width=w, height=h)
        except Exception:
            return None

    # stretch: force into bbox
    if fit_n == "stretch":
        return _add_picture(l=left, t=top, w=box_w, h=box_h)

    img_ratio = iw / ih
    box_ratio = box_w_pt / box_h_pt if box_h_pt > 0 else img_ratio

    # contain: preserve aspect, letterbox
    if fit_n == "contain":
        scale = min(box_w_pt / iw, box_h_pt / ih)
        # Heuristic: avoid upscaling small icons/logos too aggressively.
        # This keeps perceived size correct when the bbox is a loose placeholder.
        try:
            is_iconish = False
            if isinstance(img_spec, dict):
                fn = str(img_spec.get("filename") or img_spec.get("name") or "").lower()
                if any(k in fn for k in ("icon", "ico", "logo", "glyph")):
                    is_iconish = True
            # Small-ish boxes are typically icons
            if max(box_w_pt, box_h_pt) <= 160.0:
                is_iconish = True
            if is_iconish:
                scale = min(scale, 1.0)
        except Exception:
            pass
        w_pt = iw * scale
        h_pt = ih * scale
        l2 = Pt(left_pt + (box_w_pt - w_pt) / 2.0)
        t2 = Pt(top_pt + (box_h_pt - h_pt) / 2.0)
        return _add_picture(l=l2, t=t2, w=Pt(w_pt), h=Pt(h_pt))

    # cover: preserve aspect, crop
    # Place image so it fills the box and crop the overflow equally.
    pic = None
    if box_ratio >= img_ratio:
        # box is wider -> fit width, crop height
        pic = _add_picture(l=left, t=top, w=box_w, h=None)
        if pic is None:
            return None
        rendered_h_pt = box_w_pt / img_ratio
        if rendered_h_pt > 0:
            overflow_pt = max(rendered_h_pt - box_h_pt, 0.0)
            frac = overflow_pt / rendered_h_pt
            # crop fractions are 0..1
            pic.crop_top = frac / 2.0
            pic.crop_bottom = frac / 2.0
    else:
        # box is taller -> fit height, crop width
        pic = _add_picture(l=left, t=top, w=None, h=box_h)
        if pic is None:
            return None
        rendered_w_pt = box_h_pt * img_ratio
        if rendered_w_pt > 0:
            overflow_pt = max(rendered_w_pt - box_w_pt, 0.0)
            frac = overflow_pt / rendered_w_pt
            pic.crop_left = frac / 2.0
            pic.crop_right = frac / 2.0

    # Ensure positioned at the bbox top-left after fit+crop
    try:
        pic.left = left
        pic.top = top
        pic.width = box_w
        pic.height = box_h
    except Exception:
        pass

    return pic

def _load_json(path: Path) -> Any:
    return json.loads(path.read_text(encoding="utf-8"))


def _find_component_by_role(components: list[dict[str, Any]], role: str) -> dict[str, Any] | None:
    for c in components:
        if c.get("role") == role:
            return c
    return None


def _slot_defaults_for_component(comp: dict[str, Any], *, title: str, lines: list[str]) -> dict[str, Any]:
    """Fill slots for TOC/Citations components without hard-coding slot names.

    Strategy:
    - First `value_type == 'string'` slot receives `title`
    - First `value_type == 'string_list'` slot receives `lines`
    """
    out: dict[str, Any] = {}
    string_slot_name: str | None = None
    list_slot_name: str | None = None

    for s in comp.get("slots", []):
        vt = s.get("value_type")
        if vt == "string" and string_slot_name is None:
            string_slot_name = s.get("name")
        if vt == "string_list" and list_slot_name is None:
            list_slot_name = s.get("name")

    if string_slot_name is not None:
        out[string_slot_name] = title
    if list_slot_name is not None:
        out[list_slot_name] = lines
    return out


def _collect_citation_lines(blueprint: dict[str, Any]) -> list[str]:
    """Collect and dedupe citations into human-readable lines."""
    seen: set[tuple[str, str, str]] = set()
    lines: list[str] = []

    for slide in blueprint.get("slides", []):
        for c in slide.get("citations", []) or []:
            mark = str(c.get("mark", ""))
            page = c.get("page")
            chunk_id = str(c.get("chunk_id", ""))
            bbox = c.get("bbox")

            page_s = "-" if page is None else str(page)
            bbox_s = "" if bbox is None else f" bbox={bbox}"

            key = (mark, page_s, chunk_id)
            if key in seen:
                continue
            seen.add(key)

            if mark:
                lines.append(f"{mark} p.{page_s}{bbox_s} chunk={chunk_id}")
            else:
                lines.append(f"p.{page_s}{bbox_s} chunk={chunk_id}")

    return lines


def _render_pptx_dom(*, run_dir: Path, out_pptx: Path) -> None:
    """Extraction-driven renderer (A: visual fidelity first).

    - Reads run_dir/extraction.json (expects extended fields)
    - Reconstructs slides by drawing chunks in z-order
    - Does NOT use template/blueprint slots (schema strict remains untouched)
    """
    extraction = _load_json(run_dir / "extraction.json")

    # Best-effort: locate the original source .pptx so we can extract embedded media.
    media_by_sha: dict[str, tuple[bytes, str]] = {}
    src_pptx = None
    try:
        src_pptx = _find_source_pptx(run_dir)
        if src_pptx is not None and src_pptx.exists():
            media_by_sha = _build_pptx_media_sha_map(src_pptx)
    except Exception:
        media_by_sha = {}

    prs_src = None
    try:
        if src_pptx is not None and src_pptx.exists():
            prs_src = Presentation(str(src_pptx))
    except Exception:
        prs_src = None

    # Cache alignments resolved from source PPTX per (slide_no, shape_no)
    align_cache: dict[tuple[int, int], list[PP_ALIGN | None]] = {}
    vanchor_cache: dict[tuple[int, int], MSO_VERTICAL_ANCHOR | None] = {}

    # Secondary index: (byte_size, ext) -> [bytes]
    media_by_size_ext: dict[tuple[int, str], list[bytes]] = {}
    try:
        for _sha, (blob, ext) in (media_by_sha or {}).items():
            k = (len(blob), str(ext or "").lower())
            media_by_size_ext.setdefault(k, []).append(blob)
    except Exception:
        media_by_size_ext = {}

    try:
        print(f"[DOM] source_pptx={src_pptx} media_items={len(media_by_sha)}")
    except Exception:
        pass

    doc = extraction.get("document") or {}
    page_count = int(doc.get("page_count") or 0)
    page_meta = doc.get("page") or {}

    prs = Presentation()

    # Prefer slide size from extraction (EMU). Fallback to default pptx size.
    w_emu = page_meta.get("w_emu")
    h_emu = page_meta.get("h_emu")
    if isinstance(w_emu, (int, float)) and isinstance(h_emu, (int, float)) and w_emu > 0 and h_emu > 0:
        prs.slide_width = Pt(_emu_to_pt(w_emu))
        prs.slide_height = Pt(_emu_to_pt(h_emu))

    chunks = extraction.get("chunks", [])
    if not isinstance(chunks, list):
        chunks = []

    # Infer page_count if missing
    if page_count <= 0:
        pages = [int(ch.get("page") or 0) for ch in chunks if isinstance(ch, dict)]
        page_count = max(pages) if pages else 0

    def _chunk_z(ch: Any) -> float:
        if isinstance(ch, dict) and ch.get("z") is not None:
            try:
                return float(ch.get("z"))
            except Exception:
                return 0.0
        return 0.0

    def _chunk_reading_tie(ch: Any) -> tuple[int, int]:
        if not isinstance(ch, dict):
            return (0, 0)
        bbox = ch.get("bbox") or {}
        try:
            y = int(bbox.get("y") or 0)
        except Exception:
            y = 0
        try:
            x = int(bbox.get("x") or 0)
        except Exception:
            x = 0
        return (y, x)

    for page in range(1, page_count + 1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

        page_chunks = [ch for ch in chunks if isinstance(ch, dict) and int(ch.get("page", -1) or -1) == page]
        # Draw back-to-front
        page_chunks.sort(key=lambda ch: (_chunk_z(ch),) + _chunk_reading_tie(ch))

        # Best-effort: determine slide background color for BACKGROUND fills.
        page_bg_rgb: str | None = None
        try:
            # Find backmost (lowest z) chunk that covers (almost) the whole slide and has a solid fill.
            w_emu0 = page_meta.get("w_emu")
            h_emu0 = page_meta.get("h_emu")
            w_emu0 = int(w_emu0) if isinstance(w_emu0, (int, float)) else None
            h_emu0 = int(h_emu0) if isinstance(h_emu0, (int, float)) else None

            def _covers_page(b: Any) -> bool:
                if not isinstance(b, dict):
                    return False
                if w_emu0 is None or h_emu0 is None:
                    return False
                try:
                    x = int(b.get("x") or 0)
                    y = int(b.get("y") or 0)
                    w = int(b.get("w") or 0)
                    h = int(b.get("h") or 0)
                except Exception:
                    return False
                # allow small tolerance (EMU)
                tol = 50000
                return (abs(x) <= tol) and (abs(y) <= tol) and (abs(w - w_emu0) <= tol) and (abs(h - h_emu0) <= tol)

            bg_candidates: list[tuple[float, str]] = []
            for ch0 in page_chunks:
                if not isinstance(ch0, dict):
                    continue
                b0 = ch0.get("bbox")
                if not _covers_page(b0):
                    continue
                st0 = ch0.get("style")
                if not isinstance(st0, dict):
                    continue
                f0 = st0.get("fill")
                if not isinstance(f0, dict):
                    continue
                if str(f0.get("type") or "").lower() != "solid":
                    continue
                c0 = f0.get("color_rgb")
                if isinstance(c0, str) and c0.strip():
                    bg_candidates.append((_chunk_z(ch0), c0.strip().lstrip("#")))

            if bg_candidates:
                # backmost = lowest z
                bg_candidates.sort(key=lambda t: t[0])
                page_bg_rgb = bg_candidates[0][1]
        except Exception:
            page_bg_rgb = None

        for ch in page_chunks:
            kind = ch.get("kind")
            promoted_from_text = False
            promoted_shape_kind = None
            # Some decks encode background panels as empty textboxes (no runs) with a solid fill.
            # If extraction reports `kind=text` but there's no actual text, render it as a shape.
            if kind == "text":
                t0 = (ch.get("normalized_text") or ch.get("text") or "").strip()
                has_runs = False
                try:
                    ts0 = ch.get("text_struct") or {}
                    paras0 = ts0.get("paragraphs") if isinstance(ts0, dict) else None
                    if isinstance(paras0, list):
                        for p0 in paras0:
                            if not isinstance(p0, dict):
                                continue
                            runs0 = p0.get("runs")
                            if isinstance(runs0, list) and len(runs0) > 0:
                                has_runs = True
                                break
                except Exception:
                    pass

                sk0 = ch.get("shape_kind")
                st0 = ch.get("style") or {}
                fill0 = st0.get("fill") if isinstance(st0, dict) else None
                has_solid_fill = False
                try:
                    if (
                        isinstance(fill0, dict)
                        and str(fill0.get("type") or "").lower() == "solid"
                        and isinstance(fill0.get("color_rgb"), str)
                        and fill0.get("color_rgb")
                    ):
                        has_solid_fill = True
                except Exception:
                    pass

                has_nontrivial_alpha = False
                try:
                    a0 = None
                    if isinstance(fill0, dict):
                        a0 = fill0.get("alpha")
                        if a0 is None:
                            a0 = fill0.get("opacity")
                        if a0 is None:
                            a0 = fill0.get("opacity_pct")
                        if a0 is None and fill0.get("transparency") is not None:
                            # transparency is 0..1 where 1 is fully transparent
                            try:
                                a0 = 1.0 - float(fill0.get("transparency"))
                            except Exception:
                                a0 = None
                    if a0 is not None:
                        aa0 = float(a0)
                        # Treat anything other than fully opaque as a signal (0.2, 0.0, etc.)
                        if aa0 != 1.0:
                            has_nontrivial_alpha = True
                except Exception:
                    has_nontrivial_alpha = False

                # Promote to shape also for empty textboxes with solid fill or nontrivial alpha, even if shape_kind is missing.
                if (not t0) and (not has_runs) and (has_solid_fill or has_nontrivial_alpha):
                    kind = "shape"
                    promoted_from_text = True
                    promoted_shape_kind = sk0 or "rect"
            bbox_pt = _bbox_emu_to_pt_rect(ch.get("bbox"))
            if bbox_pt is None:
                continue
            x1, y1, x2, y2 = bbox_pt
            left = Pt(float(x1))
            top = Pt(float(y1))
            width = Pt(float(x2 - x1))
            height = Pt(float(y2 - y1))

            # Adapt style
            st = _style_from_extraction(ch.get("style"))


            # Resolve actual textbox fill from the source PPTX (authoritative) when possible.
            # This prevents extraction mistakes (e.g., SOLID white) from painting opaque boxes
            # that hide the real background shapes behind the text.
            if kind == "text" and prs_src is not None:
                try:
                    parsed = _parse_chunk_id_slide_shape(ch.get("chunk_id"))
                    if parsed is not None:
                        slide_no, shape_no = parsed
                        rs = _resolve_textbox_fill_from_src(prs_src, slide_no, shape_no)
                        # Apply resolved fill if we learned anything.
                        if rs.get("fill_none") is True:
                            st["fill_none"] = True
                            st.pop("fill_rgb", None)
                            st.pop("fill_alpha", None)
                        elif isinstance(rs.get("fill_rgb"), str) and rs.get("fill_rgb"):
                            st["fill_rgb"] = rs["fill_rgb"]
                            if isinstance(rs.get("fill_alpha"), (int, float)):
                                st["fill_alpha"] = float(rs["fill_alpha"])
                except Exception:
                    pass
            if kind == "text":
                # shape_kind is set when the original was a geometric auto-shape (round_rect, ellipse, rect …).
                # shape_kind is absent/None for plain textboxes (MSO_SHAPE_TYPE.TEXT_BOX).
                text_shape_kind = ch.get("shape_kind")
                is_geo_text = bool(text_shape_kind)

                if not is_geo_text:
                    # Plain textbox: force no fill, no line.
                    st["fill_none"] = True
                    st.pop("fill_rgb", None)
                    st.pop("fill_alpha", None)
                    st["line_none"] = True
                    st.pop("line_rgb", None)
                    st.pop("line_width_pt", None)
                    st.pop("line_alpha", None)
                    st.pop("line_from_src", None)
                else:
                    # Geometric auto-shape with text (round_rect, ellipse, rect …).
                    # Keep line from extraction (alpha=0 transparent lines already set line_none=True).
                    # Suppress "background" fill type (transparent); keep solid fills.
                    _fill_info = (ch.get("style") or {}).get("fill") or {}
                    _fill_type = str(_fill_info.get("type", "")).lower()
                    if not st.get("fill_rgb") and _fill_type in ("background", ""):
                        st["fill_none"] = True
                        st.pop("fill_rgb", None)
                        st.pop("fill_alpha", None)
                    # Guard: if no fill info was resolved (e.g. scheme/theme color without RGB),
                    # force fill_none to prevent p:style/fillRef (accent1) from leaking as a shadow.
                    if not any(k in st for k in ("fill_rgb", "fill_alpha", "fill_none")):
                        st["fill_none"] = True


            el_style = {"rotation_deg": ch.get("rotation_deg"), "style": st}

            # TEXT
            if kind == "text":
                if is_geo_text:
                    # Geometric auto-shape with text: preserves rounded/ellipse geometry + fill/line.
                    _ast = _autoshape_type_from_el({"shape_type": text_shape_kind})
                    shape = slide.shapes.add_shape(_ast, left, top, width, height)
                    # Copy corner-radius adj from source (preserves roundRect degree).
                    if prs_src is not None:
                        try:
                            _parsed_geo = _parse_chunk_id_slide_shape(ch.get("chunk_id"))
                            if _parsed_geo is not None:
                                _copy_adj_from_src(
                                    prs_src=prs_src,
                                    slide_no=_parsed_geo[0],
                                    shape_no=_parsed_geo[1],
                                    dst_shape=shape,
                                )
                        except Exception:
                            pass
                    # Copy line XML from source when instructed (preserves theme-based borders).
                    if prs_src is not None and st.get("line_from_src") is True:
                        try:
                            _parsed_geo = _parse_chunk_id_slide_shape(ch.get("chunk_id"))
                            if _parsed_geo is not None:
                                _copy_line_xml_from_src(
                                    prs_src=prs_src,
                                    slide_no=_parsed_geo[0],
                                    shape_no=_parsed_geo[1],
                                    dst_shape=shape,
                                )
                        except Exception:
                            pass
                else:
                    # Plain textbox: no geometry, default no-fill/no-line.
                    shape = slide.shapes.add_textbox(left, top, width, height)
                    _set_no_fill(shape)
                    _set_no_line(shape)
                _apply_common_style(shape, el_style)
                tf = shape.text_frame
                tf.clear()
                # Ensure wrapping within the textbox bounds.
                try:
                    tf.word_wrap = True
                except Exception:
                    pass
                try:
                    tf.auto_size = MSO_AUTO_SIZE.NONE
                except Exception:
                    pass
                # Remove default margins to match tight bbox from extraction and avoid accidental wrap.
                try:
                    tf.margin_left = 0
                    tf.margin_right = 0
                    tf.margin_top = 0
                    tf.margin_bottom = 0
                except Exception:
                    pass

                # Vertical alignment within textbox (top/middle/bottom)
                va = None
                try:
                    va = _vanchor_from_any(ch.get("vertical_align") or ch.get("v_align") or ch.get("valign"))
                except Exception:
                    va = None
                if va is None:
                    try:
                        va = _vanchor_from_any((ts.get("vertical_align") or ts.get("v_align") or ts.get("valign")) if isinstance(ts, dict) else None)
                    except Exception:
                        va = None
                if va is None and prs_src is not None:
                    parsed = _parse_chunk_id_slide_shape(ch.get("chunk_id"))
                    if parsed is not None:
                        slide_no, shape_no = parsed
                        key = (slide_no, shape_no)
                        va2 = vanchor_cache.get(key)
                        if key not in vanchor_cache:
                            va2 = _resolve_textbox_vanchor_from_src(prs_src, slide_no, shape_no)
                            vanchor_cache[key] = va2
                        if va2 is not None:
                            va = va2
                if va is not None:
                    try:
                        tf.vertical_anchor = va
                    except Exception:
                        pass

                ts = ch.get("text_struct") or {}
                paragraphs = ts.get("paragraphs") if isinstance(ts, dict) else None

                if isinstance(paragraphs, list) and paragraphs:
                    first_para = True
                    for p_spec in paragraphs:
                        if not isinstance(p_spec, dict):
                            continue
                        para = tf.paragraphs[0] if first_para else tf.add_paragraph()
                        first_para = False
                        # alignment
                        al = _align_from_any(p_spec.get("alignment") or p_spec.get("align") or p_spec.get("text_align"))
                        if al is None:
                            al = _align_from_any((ts.get("alignment") or ts.get("align") or ts.get("text_align")) if isinstance(ts, dict) else None)
                        if al is None:
                            al = _align_from_any(ch.get("alignment") or ch.get("align") or ch.get("text_align"))

                        # Authoritative fallback: read from source PPTX when extraction alignment is missing/unknown.
                        if al is None and prs_src is not None:
                            parsed = _parse_chunk_id_slide_shape(ch.get("chunk_id"))
                            if parsed is not None:
                                slide_no, shape_no = parsed
                                key = (slide_no, shape_no)
                                aligns = align_cache.get(key)
                                if aligns is None:
                                    aligns = _resolve_textbox_alignments_from_src(prs_src, slide_no, shape_no)
                                    align_cache[key] = aligns
                                # Use paragraph index if provided; otherwise fall back to first non-None.
                                idx = None
                                try:
                                    if isinstance(p_spec.get("index"), int):
                                        idx = int(p_spec.get("index"))
                                    elif isinstance(p_spec.get("p_index"), int):
                                        idx = int(p_spec.get("p_index"))
                                except Exception:
                                    idx = None

                                if aligns:
                                    if idx is not None and 0 <= idx < len(aligns) and aligns[idx] is not None:
                                        al = aligns[idx]
                                    else:
                                        for aa in aligns:
                                            if aa is not None:
                                                al = aa
                                                break

                        if al is not None:
                            try:
                                para.alignment = al
                            except Exception:
                                pass

                        runs = p_spec.get("runs")
                        if isinstance(runs, list) and runs:
                            for r_spec in runs:
                                if not isinstance(r_spec, dict):
                                    continue
                                txt = r_spec.get("text")
                                if txt is None:
                                    continue
                                run = para.add_run()
                                run.text = str(txt)
                                try:
                                    fn = r_spec.get("font_name")
                                    if isinstance(fn, str) and fn:
                                        run.font.name = fn
                                except Exception:
                                    pass
                                try:
                                    fs = r_spec.get("font_size_emu")
                                    if isinstance(fs, (int, float)) and fs > 0:
                                        run.font.size = Pt(_emu_to_pt(fs))
                                except Exception:
                                    pass
                                try:
                                    if r_spec.get("bold") is True:
                                        run.font.bold = True
                                    if r_spec.get("italic") is True:
                                        run.font.italic = True
                                    if r_spec.get("underline") is True:
                                        run.font.underline = True
                                except Exception:
                                    pass
                                try:
                                    c = r_spec.get("color_rgb")
                                    rgb = _rgb_from_any(c)
                                    if rgb is not None:
                                        run.font.color.rgb = rgb
                                except Exception:
                                    pass
                        else:
                            txt = (ch.get("normalized_text") or ch.get("text") or "").strip()
                            para.text = txt
                else:
                    # Fallback: plain text
                    txt = (ch.get("normalized_text") or ch.get("text") or "").strip()
                    if "\n" in txt:
                        tf.clear()
                        # keep wrapping enabled after clearing
                        try:
                            tf.word_wrap = True
                        except Exception:
                            pass
                        try:
                            tf.auto_size = MSO_AUTO_SIZE.NONE
                        except Exception:
                            pass
                        first = True
                        for line in [t.rstrip() for t in txt.splitlines()]:
                            if first:
                                tf.paragraphs[0].text = line
                                first = False
                            else:
                                tf.add_paragraph().text = line
                    else:
                        tf.text = txt

                continue

            # IMAGE
            if kind == "image":
                img = ch.get("image")
                if img is None:
                    # Some extractors may put the path directly on the chunk.
                    img = ch.get("path") or ch.get("relpath") or ch.get("filename")
                # Resolve image/icon bytes from source pptx media.
                if isinstance(img, dict):
                    sha = img.get("sha256")
                    ext = str(img.get("ext") or "").lower()
                    bsz = img.get("byte_size")

                    # 1) Exact sha match
                    if isinstance(sha, str) and sha and sha in media_by_sha:
                        blob, _ext = media_by_sha[sha]
                        img = blob
                    else:
                        # 2) Fallback: match by (byte_size, ext) if it uniquely identifies a blob
                        try:
                            if isinstance(bsz, int) and bsz > 0 and ext:
                                cands = media_by_size_ext.get((bsz, ext), [])
                                if len(cands) == 1:
                                    img = cands[0]
                        except Exception:
                            pass

                        # Debug (only when unresolved)
                        try:
                            if isinstance(img, dict):
                                # still unresolved
                                short = (sha[:10] + "..") if isinstance(sha, str) and len(sha) > 10 else sha
                                print(f"[DOM] unresolved image sha={short} ext={ext} byte_size={bsz}")
                        except Exception:
                            pass
                # If extractor emitted nested spec with path/bytes_b64, reuse add-picture helper.
                pic = _add_picture_with_fit(
                    slide=slide,
                    run_dir=run_dir,
                    left_pt=float(x1),
                    top_pt=float(y1),
                    box_w_pt=float(x2 - x1),
                    box_h_pt=float(y2 - y1),
                    img_spec=img,
                    fit=(
                        str(img.get("fit"))
                        if isinstance(img, dict) and img.get("fit")
                        else (
                            "stretch"
                            if (
                                float(x2 - x1) >= float(prs.slide_width.pt) * 0.95
                                and float(y2 - y1) >= float(prs.slide_height.pt) * 0.95
                            )
                            else "contain"
                        )
                    ),
                )
                if pic is not None:
                    _apply_common_style(pic, el_style)
                continue

            # TABLE
            if kind == "table":
                tbl = ch.get("table")
                if isinstance(tbl, dict):
                    try:
                        rows = int(tbl.get("rows", 0))
                        cols = int(tbl.get("cols", 0))
                        if rows > 0 and cols > 0:
                            table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
                            table = table_shape.table
                            cells = tbl.get("cells")
                            if isinstance(cells, list):
                                _render_table_cells(table, cells)
                            _apply_common_style(table_shape, el_style)
                    except Exception:
                        pass
                continue

            # SHAPE / LINE (best-effort)
            if kind in ("shape", "line"):
                try:
                    st_name = ch.get("shape_type") or ch.get("shape_kind") or "rect"
                    # Text boxes promoted to shapes should remain square-corner rectangles.
                    # Only preserve ellipse/line if the promotion explicitly indicated it.
                    if promoted_from_text:
                        if promoted_shape_kind == "ellipse":
                            st_name = "ellipse"
                        elif promoted_shape_kind == "line":
                            st_name = "line"
                        elif promoted_shape_kind in (
                            "round_rect", "roundrect", "rounded_rect",
                            "rounded_rectangle", "roundedrectangle",
                        ):
                            st_name = "round_rect"
                        else:
                            st_name = "rect"

                    st_style = _style_from_extraction(ch.get("style"))

                    # Resolve from source PPTX when extraction is missing/incorrect (esp. outlines).
                    if prs_src is not None:
                        parsed2 = _parse_chunk_id_slide_shape(ch.get("chunk_id"))
                        if parsed2 is not None:
                            s_no, sh_no = parsed2
                            rs2 = _resolve_shape_style_from_src(prs_src, s_no, sh_no)

                            # Decide whether the source indicates a *real* line.
                            src_has_line = False
                            try:
                                if rs2.get("line_none") is True:
                                    src_has_line = False
                                else:
                                    if rs2.get("line_width_pt") is not None:
                                        src_has_line = True
                                    if rs2.get("line_rgb") is not None:
                                        src_has_line = True
                                    if isinstance(rs2.get("line_alpha"), (int, float)) and float(rs2.get("line_alpha")) > 0:
                                        src_has_line = True
                            except Exception:
                                src_has_line = False

                            # Decide whether extraction has an explicit line spec.
                            ext_has_line_spec = any(k in st_style for k in ("line_rgb", "line_width_pt", "line_alpha"))

                            # Case A: extraction marked no-line, but source has a line -> revive from source.
                            if st_style.get("line_none") is True and src_has_line:
                                st_style.pop("line_none", None)
                                for k in ("line_rgb", "line_width_pt", "line_alpha"):
                                    if rs2.get(k) is not None:
                                        st_style[k] = rs2[k]
                                st_style["line_from_src"] = True

                            # Case B: extraction has no explicit line spec, but source has a line -> adopt source.
                            elif (not ext_has_line_spec) and src_has_line:
                                for k in ("line_rgb", "line_width_pt", "line_alpha"):
                                    if rs2.get(k) is not None:
                                        st_style[k] = rs2[k]
                                st_style.pop("line_none", None)
                                st_style["line_from_src"] = True

                            # Case C: otherwise, only fill missing keys (do not override explicit extraction).
                            else:
                                for k, v in rs2.items():
                                    if k not in st_style or st_style.get(k) is None:
                                        st_style[k] = v
                    if not any(k in st_style for k in ("fill_rgb", "fill_alpha", "fill_none")):
                        st_style["fill_none"] = True
                    el = {
                        "shape_type": st_name,
                        "rotation_deg": ch.get("rotation_deg"),
                        "style": st_style,
                    }
                    if kind == "line":
                        # Draw as a thin rectangle if we lack a true line primitive
                        min_thickness_pt = 1.0
                        w_pt = max(float(x2 - x1), min_thickness_pt)
                        h_pt = max(float(y2 - y1), min_thickness_pt)
                        shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, Pt(w_pt), Pt(h_pt))
                        try:
                            shp.fill.background()
                        except Exception:
                            pass
                        if prs_src is not None and st_style.get("line_from_src") is True:
                            try:
                                parsed3 = _parse_chunk_id_slide_shape(ch.get("chunk_id"))
                                if parsed3 is not None:
                                    s_no3, sh_no3 = parsed3
                                    _copy_line_xml_from_src(prs_src=prs_src, slide_no=s_no3, shape_no=sh_no3, dst_shape=shp)
                            except Exception:
                                pass
                        _apply_common_style(shp, el)
                    else:
                        st = _autoshape_type_from_el(el)
                        shp = slide.shapes.add_shape(st, left, top, width, height)
                        if prs_src is not None and st_style.get("line_from_src") is True:
                            try:
                                parsed3 = _parse_chunk_id_slide_shape(ch.get("chunk_id"))
                                if parsed3 is not None:
                                    s_no3, sh_no3 = parsed3
                                    _copy_line_xml_from_src(prs_src=prs_src, slide_no=s_no3, shape_no=sh_no3, dst_shape=shp)
                            except Exception:
                                pass
                        # Copy corner-radius adj from source (preserves roundRect degree).
                        if prs_src is not None:
                            try:
                                parsed3 = _parse_chunk_id_slide_shape(ch.get("chunk_id"))
                                if parsed3 is not None:
                                    _copy_adj_from_src(prs_src=prs_src, slide_no=parsed3[0], shape_no=parsed3[1], dst_shape=shp)
                            except Exception:
                                pass
                        _apply_common_style(shp, el)
                except Exception:
                    pass
                continue

            # FALLBACK SHAPE: some extractors may label background/autoshapes with other kinds.
            if kind not in ("text", "image", "table"):
                st_style = _style_from_extraction(ch.get("style"))
                if st_style:
                    try:
                        el = {
                            "shape_type": ch.get("shape_type") or ch.get("shape_kind") or "rect",
                            "rotation_deg": ch.get("rotation_deg"),
                            "style": st_style,
                        }
                        st = _autoshape_type_from_el(el)
                        shp = slide.shapes.add_shape(st, left, top, width, height)
                        _apply_common_style(shp, el)
                    except Exception:
                        pass
                continue

            # Unknown kinds are ignored in DOM mode.
            continue

    out_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_pptx)


def render_pptx(*, run_dir: Path, out_pptx: Path, mode: str = "template") -> None:
    """
    Minimal renderer v0.1:
    - reads run_dir/template.json and run_dir/blueprint.json
    - supports component role=body with component_id=comp_title_bullets
    - renders TITLE (string) and BULLETS (string_list)
    - always adds footer page number on all rendered slides (v0.1)
    - auto-inserts a TOC slide if template has a toc component and blueprint has toc[]
    - auto-appends a citations slide if template has a citations component and blueprint has citations
    """
    mode_n = str(mode or "template").lower().strip()
    if mode_n == "dom":
        # Visual-fidelity renderer driven by extraction.json
        return _render_pptx_dom(run_dir=run_dir, out_pptx=out_pptx)
    template = _load_json(run_dir / "template.json")
    blueprint = _load_json(run_dir / "blueprint.json")

    theme = template["theme"]
    page = theme["page"]
    width_pt = float(page["width_pt"])
    height_pt = float(page["height_pt"])

    prs = Presentation()
    prs.slide_width = Pt(width_pt)
    prs.slide_height = Pt(height_pt)

    # Build component index
    comp_by_id: dict[str, dict[str, Any]] = {c["component_id"]: c for c in template["components"]}

    footer_cfg = theme.get("footer", {}).get("page_number")

    # v0.1 policy: page numbers are required on all slides.
    # If the theme does not specify footer settings, enable by default with a sensible bbox.
    if isinstance(footer_cfg, dict):
        footer_enabled = bool(footer_cfg.get("enabled", True))
        footer_bbox = footer_cfg.get("bbox_pt")
        footer_font_pt = int(footer_cfg.get("font_pt", 10))
    else:
        footer_enabled = True
        footer_bbox = None
        footer_font_pt = 10

    if footer_bbox is None:
        # Default: bottom-right area
        margin_pt = 24.0
        box_w_pt = 72.0
        box_h_pt = 18.0
        x2 = width_pt - margin_pt
        y2 = height_pt - margin_pt
        x1 = x2 - box_w_pt
        y1 = y2 - box_h_pt
        footer_bbox = [x1, y1, x2, y2]

    slides = blueprint.get("slides", [])

    # Auto insert TOC / citations if components exist
    components_list = template.get("components", [])
    toc_comp = _find_component_by_role(components_list, "toc")
    cit_comp = _find_component_by_role(components_list, "citations")

    has_toc_slide = any(comp_by_id.get(s.get("component_id", ""), {}).get("role") == "toc" for s in slides)
    has_cit_slide = any(comp_by_id.get(s.get("component_id", ""), {}).get("role") == "citations" for s in slides)

    toc_items = blueprint.get("toc", []) or []
    citation_lines = _collect_citation_lines(blueprint)

    planned: list[dict[str, Any]] = []

    # Keep any explicit cover slides first
    for s in slides:
        role = comp_by_id.get(s.get("component_id", ""), {}).get("role")
        if role == "cover":
            planned.append(s)

    # Insert TOC if missing
    if (not has_toc_slide) and toc_comp is not None and toc_items:
        lines: list[str] = []
        for t in toc_items:
            # toc item may be string or object; keep minimal
            if isinstance(t, str):
                lines.append(t)
            elif isinstance(t, dict):
                title = str(t.get("title", ""))
                indent = int(t.get("level", 0))
                prefix = "  " * max(indent, 0)
                lines.append(f"{prefix}{title}")
            else:
                lines.append(str(t))

        # If we will append a citations slide, ensure it is listed in the TOC.
        will_append_citations = (not has_cit_slide) and cit_comp is not None and bool(citation_lines)
        if will_append_citations:
            # Avoid duplicates if the user already included it.
            if not any(str(x).strip() == "引用" for x in lines):
                lines.append("引用")

        planned.append(
            {
                "component_id": toc_comp["component_id"],
                "message": "",
                "slots": _slot_defaults_for_component(toc_comp, title="目次", lines=lines),
                "citations": [],
            }
        )

    # Add non-cover, non-citations slides as main content (including explicit toc if present)
    for s in slides:
        role = comp_by_id.get(s.get("component_id", ""), {}).get("role")
        if role in ("cover", "citations"):
            continue
        planned.append(s)

    # Append citations slide if missing and we have citations
    if (not has_cit_slide) and cit_comp is not None and citation_lines:
        planned.append(
            {
                "component_id": cit_comp["component_id"],
                "message": "",
                "slots": _slot_defaults_for_component(cit_comp, title="引用", lines=citation_lines),
                "citations": [],
            }
        )

    slides = planned

    rendered_no = 0

    for _, s in enumerate(slides, 1):
        comp_id = s["component_id"]
        comp = comp_by_id.get(comp_id)
        if comp is None:
            # Unknown component: skip (v0.1 minimal)
            continue

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        rendered_no += 1

        # Create shapes from layout_elements
        shapes_by_element_id: dict[str, Any] = {}
        layout_elements = comp.get("layout_elements", [])
        if not isinstance(layout_elements, list):
            layout_elements = []

        # Draw from back to front: lower z first.
        def _z(el: Any) -> float:
            if isinstance(el, dict) and el.get("z") is not None:
                try:
                    return float(el.get("z"))
                except Exception:
                    return 0.0
            return 0.0

        for el in sorted(layout_elements, key=_z):
            kind = el.get("kind")
            if kind is None:
                continue

            # Common geometry
            bbox = el.get("bbox_pt")
            if not (isinstance(bbox, list) and len(bbox) == 4):
                continue
            x1, y1, x2, y2 = bbox
            left = Pt(float(x1))
            top = Pt(float(y1))
            width = Pt(float(x2 - x1))
            height = Pt(float(y2 - y1))

            element_id = el.get("element_id")
            if not isinstance(element_id, str):
                continue

            # 1) Textbox
            if kind == "textbox":
                shape = slide.shapes.add_textbox(left, top, width, height)
                _apply_common_style(shape, el)
                shapes_by_element_id[element_id] = shape
                continue

            # 2) Image
            if kind == "image":
                img = el.get("image")
                fit = "stretch"
                if isinstance(img, dict):
                    fit = str(img.get("fit", fit) or fit)
                # Allow top-level override
                if el.get("fit") is not None:
                    fit = str(el.get("fit"))

                pic = _add_picture_with_fit(
                    slide=slide,
                    run_dir=run_dir,
                    left_pt=float(x1),
                    top_pt=float(y1),
                    box_w_pt=float(x2 - x1),
                    box_h_pt=float(y2 - y1),
                    img_spec=img,
                    fit=fit,
                )

                if pic is not None:
                    _apply_common_style(pic, el)
                    shapes_by_element_id[element_id] = pic
                continue

            # 3) Basic shapes
            if kind == "shape":
                try:
                    st = _autoshape_type_from_el(el)
                    shp = slide.shapes.add_shape(st, left, top, width, height)
                    _apply_common_style(shp, el)
                    shapes_by_element_id[element_id] = shp
                except Exception:
                    pass
                continue

            # 4) Line (very minimal): render as a thin rectangle if we lack a true line primitive
            if kind == "line":
                try:
                    # If bbox is degenerate, still draw a minimal thickness
                    min_thickness_pt = 1.0
                    w_pt = max(float(x2 - x1), min_thickness_pt)
                    h_pt = max(float(y2 - y1), min_thickness_pt)
                    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, Pt(w_pt), Pt(h_pt))
                    # Default fill none; line styled
                    try:
                        shp.fill.background()
                    except Exception:
                        pass
                    _apply_common_style(shp, el)
                    shapes_by_element_id[element_id] = shp
                except Exception:
                    pass
                continue

            # 5) Table (minimal): create a table if spec is present
            if kind == "table":
                tbl = el.get("table")
                if isinstance(tbl, dict):
                    try:
                        rows = int(tbl.get("rows", 0))
                        cols = int(tbl.get("cols", 0))
                        if rows > 0 and cols > 0:
                            table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
                            table = table_shape.table
                            cells = tbl.get("cells")
                            if isinstance(cells, list):
                                _render_table_cells(table, cells)
                            _apply_common_style(table_shape, el)
                            shapes_by_element_id[element_id] = table_shape
                    except Exception:
                        pass
                continue

            # Unknown kind: ignore in v0.1
            continue

        # Fill slots
        slots: dict[str, Any] = s.get("slots", {})
        for slot_def in comp.get("slots", []):
            name = slot_def["name"]
            value_type = slot_def["value_type"]
            element_ref = slot_def["element_ref"]
            shape = shapes_by_element_id.get(element_ref)
            if shape is None:
                continue

            tf = shape.text_frame
            tf.clear()

            val = slots.get(name)
            if val is None:
                continue

            if value_type == "string":
                p = tf.paragraphs[0]
                p.text = str(val)
            elif value_type == "string_list":
                # bullets
                for idx, item in enumerate(val):
                    if idx == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = str(item)
                    p.level = 0
            else:
                # unknown value type: ignore
                continue

            # minimal font sizing: use min_font_pt if provided
            constraints = slot_def.get("constraints", {})
            min_font_pt = constraints.get("min_font_pt")
            if isinstance(min_font_pt, int):
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(min_font_pt)

        # Footer page number
        if footer_enabled and isinstance(footer_bbox, list) and len(footer_bbox) == 4:
            x1, y1, x2, y2 = footer_bbox
            left = Pt(float(x1))
            top = Pt(float(y1))
            width = Pt(float(x2 - x1))
            height = Pt(float(y2 - y1))
            f = slide.shapes.add_textbox(left, top, width, height).text_frame
            f.text = str(rendered_no)
            for run in f.paragraphs[0].runs:
                run.font.size = Pt(footer_font_pt)

    out_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_pptx)
